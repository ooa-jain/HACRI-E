"""
What the Deeksharambh deck is allowed to claim.

Each test here stands for a figure the printed 2026 deck got wrong, or a piece
of the report it left out. They are written against the numbers, not the
layout: a slide can be redesigned freely, but it may not go back to dividing
by the wrong denominator or quietly dropping the departments at the bottom of
the table.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
import io

from app.orientation_analysis import MIN_REPORTABLE, summarize_orientation
from app.orientation_data import (
    build_report, campus_split, department_rows, redact_small_cells,
)
from app.orientation_ppt import _spread, generate_orientation_ppt


def answer(**data):
    return {"location": "📍 Bangalore", **data}


def student(dept, campus="Bangalore", level="ug", **data):
    return {"program": dept, "campus": campus, "ug_or_pg": level,
            "orientation_at": "23 Aug 2026 10:00", "orientation_at_iso": "2026-08-23",
            "name": "", "email": "",
            "data": {"location": f"📍 {campus}", **data}}


def waiting(dept, campus="Bangalore", level="ug"):
    return {"program": dept, "campus": campus, "ug_or_pg": level,
            "orientation_at": "", "orientation_at_iso": "", "name": "", "email": ""}


def deck_text(report, departments, campuses=None) -> str:
    raw = generate_orientation_ppt(
        campus="All campuses", scope="All departments", report=report,
        departments=departments, campuses=campuses or [])
    prs = Presentation(io.BytesIO(raw))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    out.append(" | ".join(c.text for c in row.cells))
            elif shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


# ── Coverage of every department ─────────────────────────────────────────────

def test_the_scoreboard_lists_every_department_not_just_the_best_twelve():
    """The 2026 deck sliced the scoreboard to twelve rows after sorting by
    vibe, so the four departments it dropped were by construction the
    lowest-rated ones."""
    filled, pending = [], []
    for i in range(16):
        dept = f"Department {i:02d}"
        filled += [student(dept, q2=10 - i * 0.4) for _ in range(MIN_REPORTABLE)]
        pending += [waiting(dept) for _ in range(5)]

    report = build_report(filled, pending)
    text = deck_text(report, redact_small_cells(department_rows(filled, pending)))

    for i in range(16):
        assert f"Department {i:02d}" in text, "a department answered and is on no slide"
    assert "Department Scoreboard (1 of 2)" in text


def test_the_scoreboard_says_how_many_departments_it_is_showing():
    filled = [student("Department of Law", q2=8) for _ in range(MIN_REPORTABLE)]
    report = build_report(filled, [])
    text = deck_text(report, redact_small_cells(department_rows(filled, [])))
    assert "All 1 departments with a student in scope" in text


# ── Small cells ──────────────────────────────────────────────────────────────

def test_a_department_below_the_threshold_keeps_its_count_and_loses_its_scores():
    filled = ([student("Big", q2=8, q29=7, q34=9) for _ in range(MIN_REPORTABLE)]
              + [student("Tiny", q2=10, q29=10, q34=10)])
    rows = {r["dept"]: r for r in redact_small_cells(
        department_rows(filled, [waiting("Tiny")]))}

    assert rows["Big"]["vibe"] == 8.0
    assert rows["Tiny"]["filled"] == 1 and rows["Tiny"]["eligible"] == 2
    for key in ("vibe", "belonging", "nps", "promoters", "detractors"):
        assert rows["Tiny"][key] is None, f"{key} still identifies one student"


def test_the_deck_explains_a_withheld_row_rather_than_printing_a_bare_dash():
    filled = ([student("Big", q2=8) for _ in range(MIN_REPORTABLE)]
              + [student("Tiny", q2=10)])
    report = build_report(filled, [])
    text = deck_text(report, redact_small_cells(department_rows(filled, [])))
    assert f"fewer than {MIN_REPORTABLE} students" in text


def test_no_department_of_one_is_named_as_the_highest_or_lowest_rated():
    """"Sports Education rated the week highest at 10.0/10" was one student."""
    rows = redact_small_cells(department_rows(
        [student("Big", q2=8) for _ in range(MIN_REPORTABLE)]
        + [student("Bigger", q2=7) for _ in range(MIN_REPORTABLE)]
        + [student("Tiny", q2=10)], []))

    high, low = _spread(rows, "vibe")
    assert {high["dept"], low["dept"]} == {"Big", "Bigger"}


# ── Denominators ─────────────────────────────────────────────────────────────

def test_a_session_share_is_reported_against_the_students_who_answered_it():
    """369 students · 35% was 35% of the 1,045 who named a session, printed
    beside a header saying 1,235 responded."""
    filled = ([student("D", q2=8, q11=["🚶 Campus Tour"]) for _ in range(7)]
              + [student("D", q2=8) for _ in range(3)])
    report = build_report(filled, [])

    assert report["count"] == 10
    assert report["highlights_answered"]["impactful"] == 7
    assert report["highlights"]["impactful"][0]["pct"] == 100.0

    text = deck_text(report, redact_small_cells(department_rows(filled, [])))
    assert "of the 7 who answered" in text


def test_the_department_count_agrees_with_the_campus_split():
    """Slide 2 printed "across 16 departments" beside campus lines adding to
    17: one counted the unmatched-reply bucket, the other did not."""
    filled = [student("Department of Law", "Bangalore", q2=8),
              student("Department of Law", "Kochi", q2=8),
              student("Department of Commerce", "Bangalore", q2=8),
              student("—", "Bangalore", q2=8)]
    report = build_report(filled, [])

    assert report["department_count"] == 2, "the unmatched bucket is not a department"
    assert report["unmatched"] == 1
    assert len({r["dept"] for r in report["departments"]}) == 3   # it is still listed

    text = deck_text(report, redact_small_cells(department_rows(filled, [])),
                     campus_split(filled, []))
    assert "across 2 departments" in text
    assert "1 replies could not be matched" in text


# ── Feedback ─────────────────────────────────────────────────────────────────

@pytest.mark.parametrize("key,question,label", [
    ("keep",             "q37", "Keep next year"),
    ("stop",             "q38", "Stop next year"),
    ("introduce",        "q39", "Introduce next year"),
    ("challenges",       "q7",  "Challenges settling in"),
    ("stressors",        "q28", "Biggest stressors"),
    ("least_connecting", "q5b", "Sessions that felt least connecting"),
    ("reasons",          "q36", "Why they scored us that way"),
])
def test_every_feedback_question_reaches_the_report_and_the_deck(key, question, label):
    """The dashboard has always drawn six feedback panels. The deck printed
    four — "stop" and "introduce" were computed and thrown away — and three
    more questions were never highlights at all."""
    filled = [student("D", q2=8, **{question: ["A distinctive answer"]})
              for _ in range(MIN_REPORTABLE)]
    report = build_report(filled, [])

    assert report["highlights"][key][0]["label"] == "A distinctive answer"
    assert report["highlights_answered"][key] == MIN_REPORTABLE

    text = deck_text(report, redact_small_cells(department_rows(filled, [])))
    assert "A distinctive answer" in text, f"{label} never reaches a slide"
    assert label in text


def test_the_feedback_slides_come_before_the_scores():
    """"Coverage wrt feedback needs to be included first." """
    filled = [student("D", q2=8, q37=["🌉 Bridge course sessions"])
              for _ in range(MIN_REPORTABLE)]
    raw = generate_orientation_ppt(
        campus="All campuses", scope="All departments",
        report=build_report(filled, []),
        departments=redact_small_cells(department_rows(filled, [])), campuses=[])

    titles = []
    for slide in Presentation(io.BytesIO(raw)).slides:
        titles += [sh.text_frame.text for sh in slide.shapes
                   if sh.has_text_frame and sh.text_frame.text.strip()]
    joined = "\n".join(titles)
    assert joined.index("What students asked us to change") < \
           joined.index("Section I — PROGRAM EFFECTIVENESS")


def test_the_deck_says_the_form_has_no_free_text_box():
    """Every orientation question is a fixed list of options. A reader asking
    "what did students write?" deserves to be told nobody could."""
    filled = [student("D", q2=8, q37=["🌉 Bridge course sessions"])
              for _ in range(MIN_REPORTABLE)]
    text = deck_text(build_report(filled, []),
                     redact_small_cells(department_rows(filled, [])))
    assert "no free-text box" in text
