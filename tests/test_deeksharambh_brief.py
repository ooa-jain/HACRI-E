"""
The Student Experience & Orientation Impact Analysis brief — a deck, a Word
document and a workbook, built to the shape of the report the Office of
Academics already circulates by hand.

Six students, by hand, so every number in the journey funnel is known:

  Asha    registered, baseline done, Deeksharambh done, post done   (full journey)
  Bilal   registered, baseline done, Deeksharambh done, post done   (full journey)
  Chitra  registered, baseline done, Deeksharambh done, post NOT done
  Dev     registered, baseline done, Deeksharambh NOT done
  Esha    registered, baseline NOT done at all
  Farah   registered, baseline done, Deeksharambh done, post done, Kochi

registered=6, baseline=5 (everyone but Esha), Deeksharambh=4, post=3,
full journey=3, "baseline but not Deeksharambh"=1 (Dev). Esha counts nowhere
in the department table either — a department's eligible pool is the
students it has actually heard from (filled or, having done the baseline,
still owed), and Esha has done neither.
"""

from __future__ import annotations

from datetime import datetime, timedelta, timezone
from typing import AsyncIterator

import pytest
import pytest_asyncio
from httpx import ASGITransport, AsyncClient
from mongomock_motor import AsyncMongoMockClient

from app import db


@pytest_asyncio.fixture
async def app_with_mock():
    db._set_client_for_tests(AsyncMongoMockClient())
    try:
        from app.main import app
        await db.init_indexes(allow_duplicate_email=True)
        yield app
    finally:
        db._reset_clients_for_tests()


@pytest_asyncio.fixture
async def admin_client(app_with_mock) -> AsyncIterator[AsyncClient]:
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        ac.cookies.set("survey_admin_session", "1")
        yield ac


@pytest_asyncio.fixture
async def client(app_with_mock) -> AsyncIterator[AsyncClient]:
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        yield ac


def _ori_answers(campus: str) -> dict:
    return {
        "location": f"📍 {campus}",
        "q2": 9,
        "q3": "🤗 Absolutely yes!",
        "q5": "🚀 Super easy",
        "q31": "✅ Yes — I know exactly who to reach",
        "q33": ["📚 Strong academic foundation & clarity", "💼 Career support & placement prep"],
        "q11": ["🏛️ University Overview & Vision"],
        "q12": ["🚌 Campus Tour"],
    }


async def _seed() -> None:
    now = datetime.now(timezone.utc)
    people = [
        ("asha@x.com", "Asha", "Department of Law", "Bangalore", db.STATUS_POST_DONE, True, True),
        ("bilal@x.com", "Bilal", "Department of Commerce", "Bangalore", db.STATUS_POST_DONE, True, True),
        ("chitra@x.com", "Chitra", "Department of Law", "Bangalore", db.STATUS_PRE_DONE, True, False),
        ("dev@x.com", "Dev", "Department of Commerce", "Bangalore", db.STATUS_PRE_DONE, False, False),
        ("esha@x.com", "Esha", "Department of Law", "Bangalore", None, False, False),
        ("farah@x.com", "Farah", "Department of Law", "Kochi", db.STATUS_POST_DONE, True, True),
    ]
    for i, (email, name, program, campus, status, orientation, post) in enumerate(people):
        await db.get_db()["users"].insert_one({
            "email": email, "name": name, "program": program, "ug_or_pg": "ug",
            "location": campus, "status": status, "created_at": now,
            "pre_submitted_at": now if status else None,
        })
        if status in (db.STATUS_PRE_DONE, db.STATUS_POST_DONE):
            await db.get_db()["pre_responses"].insert_one(
                {"email": email, "submitted_at": now, "fields": {}})
        if status == db.STATUS_POST_DONE:
            await db.get_db()["post_responses"].insert_one(
                {"email": email, "submitted_at": now, "fields": {}})
        if orientation:
            await db.get_db()["orientation_responses"].insert_one({
                "email": email, "name": name, "submitted_at": now - timedelta(hours=i),
                "data": _ori_answers(campus),
            })


# ── The aggregator ────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_deeksharambh_brief_journey_reconciles(app_with_mock):
    await _seed()
    from app.orientation_data import deeksharambh_brief

    data = await deeksharambh_brief(campus="")
    stages = {s["key"]: s["count"] for s in data["journey"]["stages"]}
    assert stages == {"registered": 6, "pre": 5, "orientation": 4, "post": 3}
    assert data["full_journey"] == 3            # Asha, Bilal, Farah
    assert data["pre_not_deeksharambh"] == 1     # Dev
    assert data["deeksharambh_count"] == 4
    assert data["pre_completed"] == 5
    assert data["deeksharambh_pct_of_pre"] == 80

    # All four Deeksharambh answers agree, so the top option is unambiguous.
    assert data["welcomed"]["label"] == "🤗 Absolutely yes!"
    assert data["welcomed"]["pct"] == 100.0
    assert data["high_vibe_pct"] == 100.0  # everyone rated 9

    depts = {d["dept"]: d for d in data["departments"]}
    assert depts["Department of Law"]["eligible"] == 3     # Asha, Chitra, Farah — not Esha
    assert depts["Department of Law"]["filled"] == 3
    assert depts["Department of Commerce"]["eligible"] == 2  # Bilal, Dev
    assert depts["Department of Commerce"]["filled"] == 1    # Bilal only, not Dev

    labels = {o["label"] for o in data["expectations"]}
    assert "📚 Strong academic foundation & clarity" in labels


@pytest.mark.asyncio
async def test_deeksharambh_brief_scoped_to_one_campus(app_with_mock):
    await _seed()
    from app.orientation_data import deeksharambh_brief

    data = await deeksharambh_brief(campus="Kochi")
    stages = {s["key"]: s["count"] for s in data["journey"]["stages"]}
    assert stages == {"registered": 1, "pre": 1, "orientation": 1, "post": 1}
    assert [d["dept"] for d in data["departments"]] == ["Department of Law"]


# ── The exports ───────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_brief_ppt_has_the_seven_slides_with_real_numbers(admin_client):
    await _seed()
    r = await admin_client.get("/admin/survey/deeksharambh-brief-ppt")
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    assert "Student_Experience_Brief" in r.headers["content-disposition"]

    from io import BytesIO
    from pptx import Presentation

    deck = Presentation(BytesIO(r.content))
    assert len(deck.slides) == 7  # cover, journey, departments, theme, overall, expectations, conclusion

    words = " ".join(
        shape.text_frame.text
        for slide in deck.slides for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Human + AI: Building Your Future" in words
    assert "STUDENT JOURNEY" in words
    assert "6" in words              # registered
    assert "Department of Law" in words
    assert "The Vibe Check" in words  # a real section title, not invented
    assert "THANK YOU" in words


@pytest.mark.asyncio
async def test_brief_docx_carries_the_same_journey_numbers(admin_client):
    await _seed()
    r = await admin_client.get("/admin/survey/deeksharambh-brief-docx")
    assert r.status_code == 200
    assert "wordprocessingml" in r.headers["content-type"]

    from io import BytesIO
    from docx import Document

    doc = Document(BytesIO(r.content))
    journey_table = doc.tables[0]
    header_cells = [c.text for c in journey_table.rows[0].cells]
    assert "Registered\n6" in header_cells
    assert "Deeksharambh\n4" in header_cells


@pytest.mark.asyncio
async def test_brief_xlsx_has_one_sheet_per_section(admin_client):
    await _seed()
    r = await admin_client.get("/admin/survey/deeksharambh-brief-xlsx")
    assert r.status_code == 200
    assert "spreadsheetml" in r.headers["content-type"]

    from io import BytesIO
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(r.content))
    assert wb.sheetnames == ["Overview", "Departments", "What students said", "Expectations"]
    dept_sheet = wb["Departments"]
    assert dept_sheet["A1"].value == "Department"
    depts = {row[0].value: row[1].value for row in dept_sheet.iter_rows(min_row=2)}
    assert depts["Department of Law"] == 3


@pytest.mark.asyncio
async def test_brief_downloads_need_the_admin_cookie(app_with_mock):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        for path in ("deeksharambh-brief-ppt", "deeksharambh-brief-docx", "deeksharambh-brief-xlsx"):
            assert (await anon.get(f"/admin/survey/{path}")).status_code == 403


# ── Reachable from the impact page ───────────────────────────────────────

@pytest.mark.asyncio
async def test_the_impact_page_links_to_the_brief(client):
    await _seed()
    from app.routes.shared_analysis import get_vibe_token

    await _seed_matched_for_impact()
    token = get_vibe_token("")
    page = await client.get("/shared/impact", params={"campus": "", "token": token})
    assert page.status_code == 200
    assert "Human + AI: Building Your Future" in page.text
    assert "/shared/impact/brief-ppt?campus=&amp;token=" in page.text

    r = await client.get("/shared/impact/brief-ppt", params={"campus": "", "token": token})
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]

    for path in ("brief-ppt", "brief-docx", "brief-xlsx"):
        assert (await client.get(f"/shared/impact/{path}",
                                 params={"campus": "", "token": "wrong"})).status_code == 403


async def _seed_matched_for_impact() -> None:
    """The impact page needs at least one pre+post match to render at all."""
    now = datetime.now(timezone.utc)
    await db.get_db()["users"].insert_one({
        "email": "matched@x.com", "name": "Matched", "program": "Department of Law",
        "ug_or_pg": "ug", "location": "Bangalore", "status": db.STATUS_POST_DONE, "created_at": now,
    })
    from app.hacri_e2_compat import SCHEMA
    await db.get_db()["pre_responses"].insert_one(
        {"email": "matched@x.com", "submitted_at": now, "fields": {k: 3 for k in SCHEMA}})
    await db.get_db()["post_responses"].insert_one(
        {"email": "matched@x.com", "submitted_at": now, "fields": {k: 4 for k in SCHEMA}})
