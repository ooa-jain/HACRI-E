"""
The plain department headcount — a slide deck and a Word document, with no
score in either. Built from the same seed test_orientation_admin.py uses, so
the totals here are known by hand:

  Bangalore: Asha (Law, filled) · Bilal (Commerce, filled) · Dev (pending)
  Kochi:     Chitra (Law, filled)
  Esha never finished the baseline, so she is in neither cohort.

registered = 4, answered = 3, pending = 1; Bangalore has 2 departments of one
student each, Kochi has one.
"""

from __future__ import annotations

from datetime import datetime, timedelta, timezone
from typing import AsyncIterator

import pytest
import pytest_asyncio
from httpx import ASGITransport, AsyncClient
from mongomock_motor import AsyncMongoMockClient

from app import db


@pytest_asyncio.fixture
async def app_with_mock():
    db._set_client_for_tests(AsyncMongoMockClient())
    try:
        from app.main import app
        await db.init_indexes(allow_duplicate_email=True)
        yield app
    finally:
        db._reset_clients_for_tests()


@pytest_asyncio.fixture
async def client(app_with_mock) -> AsyncIterator[AsyncClient]:
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        ac.cookies.set("survey_admin_session", "1")
        yield ac


async def _seed() -> None:
    now = datetime.now(timezone.utc)
    people = [
        ("a@x.com", "Asha", "Department of Law", "Bangalore", db.STATUS_PRE_DONE, True),
        ("b@x.com", "Bilal", "Department of Commerce", "Bangalore", db.STATUS_POST_DONE, True),
        ("c@x.com", "Chitra", "Department of Law", "Kochi", db.STATUS_PRE_DONE, True),
        ("d@x.com", "Dev", "Department of Law", "Bangalore", db.STATUS_PRE_DONE, False),
        ("e@x.com", "Esha", "Department of Law", "Kochi", None, False),
    ]
    for i, (email, name, program, campus, status, filled) in enumerate(people):
        await db.get_db()["users"].insert_one({
            "email": email, "name": name, "program": program, "ug_or_pg": "ug",
            "location": campus, "status": status, "created_at": now,
            "pre_submitted_at": now if status else None,
            "orientation_submitted": filled,
        })
        if filled:
            await db.get_db()["orientation_responses"].insert_one({
                "email": email, "name": name,
                "submitted_at": now - timedelta(hours=i),
                "data": {"location": f"📍 {campus}", "q2": 8},
            })


# ── The aggregator ────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_department_count_summary_reconciles(app_with_mock):
    await _seed()
    from app.orientation_data import department_count_summary

    data = await department_count_summary(campus="")
    assert (data["total_registered"], data["total_answered"], data["total_pending"]) == (4, 3, 1)
    assert data["response_rate"] == 75.0
    assert (data["ug"], data["pg"]) == (3, 0)

    by_campus = {row["campus"]: row for row in data["campuses"]}
    assert set(by_campus) == {"Bangalore", "Kochi"}
    assert by_campus["Bangalore"]["count"] == 2
    assert by_campus["Kochi"]["count"] == 1
    assert {d["dept"]: d["count"] for d in by_campus["Bangalore"]["departments"]} == {
        "Department of Law": 1, "Department of Commerce": 1,
    }
    assert {d["dept"]: d["count"] for d in by_campus["Kochi"]["departments"]} == {
        "Department of Law": 1,
    }
    # The combined total, campus-blind: two students named "Department of Law".
    assert {d["dept"]: d["count"] for d in data["departments"]} == {
        "Department of Law": 2, "Department of Commerce": 1,
    }


@pytest.mark.asyncio
async def test_department_count_summary_scoped_to_one_campus(app_with_mock):
    await _seed()
    from app.orientation_data import department_count_summary

    data = await department_count_summary(campus="Kochi")
    assert (data["total_answered"], data["total_pending"]) == (1, 0)
    assert [row["campus"] for row in data["campuses"]] == ["Kochi"]
    assert data["campuses"][0]["departments"] == [{"dept": "Department of Law", "count": 1}]


# ── The PPTX ──────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_count_ppt_is_a_real_deck_with_no_scores(client):
    await _seed()
    r = await client.get("/admin/survey/orientation-count-ppt")
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    assert "Department_Count" in r.headers["content-disposition"]

    from io import BytesIO
    from pptx import Presentation

    deck = Presentation(BytesIO(r.content))
    assert len(deck.slides) == 4  # cover, overview, Bangalore, Kochi

    words = " ".join(
        shape.text_frame.text
        for slide in deck.slides for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "DEPARTMENT" in words and "RESPONSE COUNT" in words
    assert "Bangalore" in words and "Kochi" in words
    # No score of any kind belongs in a count-only deck.
    for banned in ("Vibe", "NPS", "Net Promoter", "Belonging"):
        assert banned not in words

    tables = [shape.table for slide in deck.slides for shape in slide.shapes if shape.has_table]
    assert len(tables) == 2  # one department table per campus
    header = [cell.text for cell in tables[0].rows[0].cells]
    assert header == ["#", "Department", "Responses"]
    table_words = " ".join(cell.text for t in tables for row in t.rows for cell in row.cells)
    assert "Department of Law" in table_words
    assert "Department of Commerce" in table_words


@pytest.mark.asyncio
async def test_count_ppt_can_be_scoped_to_one_campus(client):
    await _seed()
    r = await client.get("/admin/survey/orientation-count-ppt", params={"campus": "Kochi"})
    assert r.status_code == 200
    assert "Kochi" in r.headers["content-disposition"]

    from io import BytesIO
    from pptx import Presentation

    deck = Presentation(BytesIO(r.content))
    assert len(deck.slides) == 3  # cover, overview, one campus table


# ── The DOCX ──────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_count_docx_carries_the_same_numbers_as_the_deck(client):
    await _seed()
    r = await client.get("/admin/survey/orientation-count-docx")
    assert r.status_code == 200
    assert "wordprocessingml" in r.headers["content-type"]
    assert "Department_Count" in r.headers["content-disposition"]

    from io import BytesIO
    from docx import Document

    doc = Document(BytesIO(r.content))
    assert len(doc.tables) == 4  # cover, overview, Bangalore, Kochi

    cover_text = doc.tables[0].rows[0].cells[0].text
    assert "DEPARTMENT RESPONSE COUNT" in cover_text

    overview = {row.cells[0].text: row.cells[1].text for row in doc.tables[1].rows}
    assert overview["Registered"] == "4"
    assert overview["Answered"] == "3"
    assert overview["Still pending"] == "1"

    bang_rows = {row.cells[1].text: row.cells[2].text for row in doc.tables[2].rows[1:]}
    assert bang_rows["Department of Law"] == "1"
    assert bang_rows["Department of Commerce"] == "1"
    assert bang_rows["Total"] == "2"


@pytest.mark.asyncio
async def test_count_downloads_need_the_admin_cookie(app_with_mock):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        assert (await anon.get("/admin/survey/orientation-count-ppt")).status_code == 403
        assert (await anon.get("/admin/survey/orientation-count-docx")).status_code == 403
