"""
The outcome and impact report — the journey through the three surveys, the
baseline the cohort started from, and what changed after the workshop.
"""

from __future__ import annotations

from datetime import datetime, timezone
from typing import AsyncIterator

import pytest
import pytest_asyncio
from httpx import ASGITransport, AsyncClient
from mongomock_motor import AsyncMongoMockClient

from app import db
from app.cohort_analysis import build_cohort_report, journey_block, movement_block, score_block, scores_of
from app.routes.shared_analysis import get_cohort_token


@pytest_asyncio.fixture
async def app_with_mock():
    db._set_client_for_tests(AsyncMongoMockClient())
    try:
        from app.main import app
        await db.init_indexes(allow_duplicate_email=True)
        yield app
    finally:
        db._reset_clients_for_tests()


@pytest_asyncio.fixture
async def client(app_with_mock) -> AsyncIterator[AsyncClient]:
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        ac.cookies.set("survey_admin_session", "1")
        yield ac


def _sheet(value: int) -> dict:
    """A full Likert sheet that scores exactly `value` on both dimensions.

    Reversed items score 6 − raw, so they have to be filled in mirrored or the
    average lands somewhere in the middle.
    """
    from app.hacri_e2_compat import SCHEMA
    return {key: str(6 - value if reversed_item else value)
            for key, (_dim, reversed_item) in SCHEMA.items()}


async def _seed() -> None:
    """Six students at two campuses, at every stage of the journey."""
    now = datetime.now(timezone.utc)
    # (email, dept, campus, pre?, orientation?, post?)
    people = [
        ("a@x.com", "Department of Law", "Bangalore", True, True, True),
        ("b@x.com", "Department of Law", "Bangalore", True, True, True),
        ("c@x.com", "Department of Law", "Bangalore", True, True, False),
        ("d@x.com", "Department of Commerce", "Bangalore", True, False, False),
        ("e@x.com", "Department of Commerce", "Bangalore", False, False, False),
        ("f@x.com", "Department of Law", "Kochi", True, True, True),
    ]
    for email, dept, campus, pre, orientation, post in people:
        status = (db.STATUS_POST_DONE if post else
                  db.STATUS_PRE_DONE if pre else None)
        await db.get_db()["users"].insert_one({
            "email": email, "name": email[0].upper(), "program": dept,
            "ug_or_pg": "ug", "location": campus, "status": status,
            "created_at": now, "pre_submitted_at": now if pre else None,
            "post_submitted_at": now if post else None,
            "orientation_submitted": orientation,
        })
        if pre:
            await db.get_db()["pre_responses"].insert_one({
                "email": email, "submitted_at": now, "fields": _sheet(2)})
        if orientation:
            await db.get_db()["orientation_responses"].insert_one({
                "email": email, "submitted_at": now, "data": {"q2": 8}})
        if post:
            await db.get_db()["post_responses"].insert_one({
                "email": email, "submitted_at": now, "fields": _sheet(4)})


# ── The maths ────────────────────────────────────────────────────────────────
def _rows(spec: list[tuple[str, bool, bool, bool]]) -> list[dict]:
    rows = []
    for dept, pre, orientation, post in spec:
        rows.append({
            "email": f"{len(rows)}@x.com", "program": dept, "campus": "Bangalore",
            "ug_or_pg": "ug",
            "status": (db.STATUS_POST_DONE if post else
                       db.STATUS_PRE_DONE if pre else "not_started"),
            "pre": _sheet(2) if pre else None,
            "post": _sheet(4) if post else None,
            "orientation": orientation,
        })
    return rows


def test_journey_counts_each_step_and_what_is_still_owed():
    rows = _rows([("Law", True, True, True), ("Law", True, True, False),
                  ("Law", True, False, False), ("Law", False, False, False)])
    j = journey_block(rows)

    assert [s["count"] for s in j["stages"]] == [4, 3, 2, 1]
    assert [s["label"] for s in j["stages"]] == \
        ["Registered", "Baseline", "Deeksharambh", "Post survey"]
    # Each step is measured against the one before it, not against registration.
    assert j["stages"][2]["of_previous"] == round(100 * 2 / 3, 1)
    assert j["stages"][2]["of_registered"] == 50.0
    assert j["stages"][3]["dropped"] == 1
    assert (j["pending_pre"], j["pending_orientation"], j["pending_post"]) == (1, 1, 2)
    assert j["completion_pct"] == 25.0


def test_journey_of_an_empty_scope_is_all_zeroes():
    j = journey_block([])
    assert j["registered"] == 0
    assert j["completion_pct"] == 0.0
    assert [s["count"] for s in j["stages"]] == [0, 0, 0, 0]


def test_score_block_bins_the_scores_and_names_the_quadrant():
    rows = _rows([("Law", True, True, True), ("Law", True, True, True)])
    block = score_block(scores_of(rows, "pre"), "Baseline")

    assert block["scored"] == 2
    assert block["avg_lit"] == 2.0 and block["avg_read"] == 2.0
    filled = [b for b in block["literacy"] if b["count"]]
    assert [(b["label"], b["count"]) for b in filled] == [("2.0–2.5", 2)]
    assert [q for q in block["quadrants"] if q["count"]][0]["label"] == "Q3: AI Novice"
    assert [b for b in block["bands"] if b["count"]][0]["label"] == "AI-Anxious Student"


def test_movement_only_counts_students_who_did_both():
    rows = _rows([("Law", True, True, True), ("Law", True, True, False)])
    m = movement_block(rows)

    assert m["matched"] == 1
    assert m["gained"] == 1 and m["declined"] == 0
    assert m["avg_delta_lit"] == 2.0        # sheet of 2s to a sheet of 4s
    assert m["moved_quadrant"] == 1         # Novice to Champion


def test_report_ranks_departments_and_names_the_leaders():
    rows = _rows([("Law", True, True, True), ("Law", True, True, True),
                  ("Commerce", True, False, False), ("Commerce", False, False, False)])
    report = build_cohort_report(rows, campus="Bangalore")

    assert [d["dept"] for d in report["departments"]] == ["Law", "Commerce"]
    law = report["departments"][0]
    assert (law["registered"], law["pre"], law["orientation"], law["post"]) == (2, 2, 2, 2)
    assert law["completion_pct"] == 100.0
    assert law["delta"] == 2.0
    assert report["leaders"]["most_complete"]["dept"] == "Law"
    assert report["leaders"]["least_complete"]["dept"] == "Commerce"
    assert report["scope"]["campus"] == "Bangalore"


def test_the_report_carries_no_student_identities():
    rows = _rows([("Law", True, True, True)])
    report = build_cohort_report(rows)
    assert "email" not in repr(report)
    assert "@x.com" not in repr(report)


# ── The endpoints ────────────────────────────────────────────────────────────
@pytest.mark.asyncio
async def test_the_admin_endpoint_scopes_by_campus(client):
    await _seed()

    everyone = (await client.get("/admin/api/cohort")).json()
    assert everyone["journey"]["registered"] == 6
    assert [s["count"] for s in everyone["journey"]["stages"]] == [6, 5, 4, 3]

    kochi = (await client.get("/admin/api/cohort", params={"campus": "Kochi"})).json()
    assert kochi["journey"]["registered"] == 1
    assert kochi["scope"]["campus"] == "Kochi"

    law = (await client.get("/admin/api/cohort",
                            params={"campus": "Bangalore", "dept": "Department of Law"})).json()
    assert law["journey"]["registered"] == 3
    assert law["outcome"]["avg_overall"] == 2.0
    assert law["impact"]["avg_overall"] == 4.0
    assert law["movement"]["matched"] == 2


@pytest.mark.asyncio
async def test_the_campus_rows_cover_both_campuses(client):
    await _seed()
    rows = {c["campus"]: c for c in (await client.get("/admin/api/cohort")).json()["campuses"]}
    assert rows["Bangalore"]["registered"] == 5
    assert rows["Kochi"]["registered"] == 1
    assert rows["Kochi"]["post"] == 1


@pytest.mark.asyncio
async def test_the_deck_downloads(client):
    await _seed()
    r = await client.get("/admin/survey/cohort-ppt", params={"campus": "Bangalore"})
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]

    from io import BytesIO
    from pptx import Presentation

    deck = Presentation(BytesIO(r.content))
    assert len(deck.slides) >= 8
    words = " ".join(shape.text_frame.text
                     for slide in deck.slides for shape in slide.shapes
                     if shape.has_text_frame)
    assert "AI Literacy Outcome and Impact" in words
    assert "Bangalore" in words


@pytest.mark.asyncio
async def test_the_deck_survives_an_empty_scope(client):
    await _seed()
    r = await client.get("/admin/survey/cohort-ppt", params={"dept": "Department of Design"})
    assert r.status_code == 200
    assert len(r.content) > 10_000


@pytest.mark.asyncio
async def test_share_links_and_the_public_page(app_with_mock):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as admin:
        admin.cookies.set("survey_admin_session", "1")
        await _seed()
        links = (await admin.get("/admin/api/cohort/share-links")).json()["links"]

    assert [row["campus"] for row in links] == ["All campuses", "Bangalore", "Kochi"]
    assert len({row["url"].split("token=")[1] for row in links}) == 3

    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        good = {"campus": "Bangalore", "token": get_cohort_token("Bangalore")}
        page = await anon.get("/shared/cohort", params=good)
        assert page.status_code == 200
        assert "Outcome and Impact" in page.text

        data = (await anon.get("/shared/cohort/data", params=good)).json()
        assert data["report"]["journey"]["registered"] == 5
        assert "Department of Law" in data["dept_options"]
        # Aggregate only — nobody is named.
        assert "@x.com" not in (await anon.get("/shared/cohort/data", params=good)).text

        deck = await anon.get("/shared/cohort/ppt", params=good)
        assert deck.status_code == 200


@pytest.mark.asyncio
@pytest.mark.parametrize("params", [
    {"campus": "Bangalore", "token": "beefbeefbeefbeef"},
    {"campus": "Bangalore", "token": ""},
    {"campus": "Kochi", "token": get_cohort_token("Bangalore")},
    # An orientation token must not open the cohort report.
    {"campus": "Bangalore", "token": "orientation"},
])
async def test_a_wrong_token_is_refused(app_with_mock, params):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        for path in ("/shared/cohort", "/shared/cohort/data", "/shared/cohort/ppt"):
            assert (await anon.get(path, params=params)).status_code == 403


@pytest.mark.asyncio
async def test_cohort_endpoints_need_the_admin_cookie(app_with_mock):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        for path in ("/admin/api/cohort", "/admin/api/cohort/share-links",
                     "/admin/survey/cohort-ppt"):
            assert (await anon.get(path)).status_code == 403


@pytest.mark.asyncio
async def test_the_renderer_and_stylesheet_are_served(client):
    for path in ("/static/js/cohort_report.js", "/static/css/cohort_report.css"):
        r = await client.get(path)
        assert r.status_code == 200
        assert len(r.content) > 1000


def test_charts_write_pngs_even_with_nothing_to_draw(tmp_path):
    from app.cohort_charts import (
        plot_campus_split, plot_departments, plot_distribution, plot_journey,
        plot_mix, plot_movement, plot_pre_post,
    )

    empty_block = score_block([], "Baseline")
    checks = [
        ("journey", plot_journey, (journey_block([]),)),
        ("dist", plot_distribution, (empty_block,)),
        ("prepost", plot_pre_post, (empty_block, empty_block)),
        ("mix", plot_mix, (empty_block, empty_block)),
        ("move", plot_movement, (movement_block([]),)),
        ("dept", plot_departments, ([],)),
        ("campus", plot_campus_split, ([],)),
    ]
    for name, fn, args in checks:
        out = fn(*args, tmp_path / f"{name}.png")
        assert out.exists() and out.stat().st_size > 0
