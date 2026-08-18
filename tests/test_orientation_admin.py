"""
Survey admin's orientation report — campus summary, analysis, who filled it,
and mailing that cohort.
"""

from __future__ import annotations

from datetime import datetime, timedelta, timezone
from typing import AsyncIterator

import pytest
import pytest_asyncio
from httpx import ASGITransport, AsyncClient
from mongomock_motor import AsyncMongoMockClient

from app import db
from app.orientation_analysis import normalize_campus, summarize_orientation


@pytest_asyncio.fixture
async def app_with_mock():
    db._set_client_for_tests(AsyncMongoMockClient())
    try:
        from app.main import app
        await db.init_indexes(allow_duplicate_email=True)
        yield app
    finally:
        db._reset_clients_for_tests()


@pytest_asyncio.fixture
async def client(app_with_mock) -> AsyncIterator[AsyncClient]:
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        ac.cookies.set("survey_admin_session", "1")
        yield ac


def _answers(*, campus: str, vibe: int, nps: int) -> dict:
    return {
        "location": f"📍 {campus}",
        "q1": ["🔥 Inspiring", "🚀 Exciting"],
        "q2": vibe,
        "q3": "🙂 Yes, mostly",
        "q8": {"Campus life": "Yes", "Academics": "Somewhat"},
        "q16": 4,
        "q29": 8,
        "q32": 9,
        "q34": nps,
        "q37": ["🎪 Student Club Fair"],
        "q41": "A JAIN Explorer",
    }


async def _seed() -> None:
    """Two Bangalore students who filled it, one Kochi who did, one who hasn't."""
    now = datetime.now(timezone.utc)
    people = [
        ("a@x.com", "Asha", "Department of Law",      "Bangalore", db.STATUS_PRE_DONE,  True,  10, 10),
        ("b@x.com", "Bilal", "Department of Commerce", "Bangalore", db.STATUS_POST_DONE, True,  6,  5),
        ("c@x.com", "Chitra", "Department of Law",     "Kochi",     db.STATUS_PRE_DONE,  True,  8,  9),
        ("d@x.com", "Dev", "Department of Law",        "Bangalore", db.STATUS_PRE_DONE,  False, 0,  0),
        # Never finished the baseline — not part of either cohort.
        ("e@x.com", "Esha", "Department of Law",       "Kochi",     None,                False, 0,  0),
    ]
    for i, (email, name, program, campus, status, filled, vibe, nps) in enumerate(people):
        await db.get_db()["users"].insert_one({
            "email": email, "name": name, "program": program, "ug_or_pg": "ug",
            "location": campus, "status": status, "created_at": now,
            "pre_submitted_at": now if status else None,
            "orientation_submitted": filled,
        })
        if filled:
            await db.get_db()["orientation_responses"].insert_one({
                "email": email, "name": name,
                "submitted_at": now - timedelta(hours=i),
                "data": _answers(campus=campus, vibe=vibe, nps=nps),
            })


# ── The analysis itself ───────────────────────────────────────────────────────
@pytest.mark.parametrize("raw, expected", [
    ("📍 Bangalore", "Bangalore"),
    ("bangalore", "Bangalore"),
    ("  Kochi ", "Kochi"),
    ("📍Kochi", "Kochi"),
    ("", ""),
    (None, ""),
    ("Mysuru", ""),
])
def test_normalize_campus(raw, expected):
    assert normalize_campus(raw) == expected


def test_summarize_counts_scales_multis_and_nps():
    records = [
        _answers(campus="Bangalore", vibe=10, nps=10),
        _answers(campus="Bangalore", vibe=6, nps=5),
    ]
    out = summarize_orientation(records)

    assert out["count"] == 2
    assert out["headline"]["vibe"] == 8.0          # (10 + 6) / 2
    assert out["headline"]["belonging"] == 8.0
    assert out["headline"]["promoters"] == 1
    assert out["headline"]["detractors"] == 1
    assert out["headline"]["nps"] == 0.0           # 50% promoters − 50% detractors

    questions = {q["key"]: q for s in out["sections"] for q in s["questions"]}
    # Multi-select: both students picked both words.
    assert questions["q1"]["answered"] == 2
    assert questions["q1"]["options"][0]["count"] == 2
    # Single choice.
    assert questions["q3"]["options"][0] == {"label": "🙂 Yes, mostly", "count": 2, "pct": 100.0}
    # Matrix keeps one distribution per statement.
    rows = {r["label"]: r for r in questions["q8"]["rows"]}
    assert rows["Campus life"]["options"][0]["label"] == "Yes"


def test_summarize_empty_cohort_is_safe():
    out = summarize_orientation([])
    assert out["count"] == 0
    assert out["sections"] == []
    assert out["headline"]["vibe"] is None
    assert out["headline"]["nps"] is None


def test_summarize_ignores_unanswered_questions():
    out = summarize_orientation([{"q2": 7}])
    keys = {q["key"] for s in out["sections"] for q in s["questions"]}
    assert keys == {"q2"}


# ── The admin endpoints ───────────────────────────────────────────────────────
@pytest.mark.asyncio
async def test_campus_summary_splits_bangalore_and_kochi(client):
    await _seed()
    r = await client.get("/admin/api/orientation/campuses")
    assert r.status_code == 200
    body = r.json()

    cards = {c["campus"]: c for c in body["campuses"]}
    assert cards["Bangalore"]["filled"] == 2
    assert cards["Bangalore"]["pending"] == 1        # Dev
    assert cards["Kochi"]["filled"] == 1
    assert cards["Kochi"]["pending"] == 0            # Esha never did the baseline
    assert body["all"]["filled"] == 3
    assert body["all"]["pending"] == 1
    assert cards["Bangalore"]["vibe"] == 8.0


@pytest.mark.asyncio
async def test_report_is_scoped_to_the_chosen_campus(client):
    await _seed()
    r = await client.get("/admin/api/orientation/report", params={"campus": "Kochi"})
    body = r.json()

    assert body["campus"] == "Kochi"
    assert body["count"] == 1
    assert body["coverage"] == {"filled": 1, "pending": 0, "eligible": 1, "pct": 100.0}
    assert body["headline"]["vibe"] == 8.0
    assert [d["dept"] for d in body["departments"]] == ["Department of Law"]


@pytest.mark.asyncio
async def test_report_follows_the_department_filter(client):
    await _seed()
    r = await client.get("/admin/api/orientation/report",
                         params={"campus": "Bangalore", "dept": "Department of Law"})
    body = r.json()
    assert body["count"] == 1                       # Asha only; Bilal is Commerce
    assert body["coverage"]["pending"] == 1         # Dev is Law too


@pytest.mark.asyncio
async def test_filled_and_pending_lists(client):
    await _seed()

    filled = (await client.get("/admin/api/orientation/students",
                               params={"campus": "Bangalore", "group": "filled"})).json()
    assert [s["name"] for s in filled["students"]] == ["Asha", "Bilal"]  # newest first
    assert filled["students"][0]["answered"] > 5
    assert filled["students"][0]["vibe"] == 10

    pending = (await client.get("/admin/api/orientation/students",
                                params={"campus": "Bangalore", "group": "pending"})).json()
    assert [s["name"] for s in pending["students"]] == ["Dev"]
    assert pending["filled"] == 2


@pytest.mark.asyncio
async def test_department_analysis_ranks_and_ignores_the_dept_filter(client):
    await _seed()
    r = await client.get("/admin/api/orientation/departments",
                         params={"campus": "Bangalore", "dept": "Department of Law"})
    body = r.json()

    # Both Bangalore departments are present despite the dept parameter — this
    # view exists to compare them.
    by_dept = {row["dept"]: row for row in body["departments"]}
    assert set(by_dept) == {"Department of Law", "Department of Commerce"}

    law = by_dept["Department of Law"]
    assert (law["filled"], law["pending"], law["eligible"]) == (1, 1, 2)
    assert law["pct"] == 50.0
    assert law["vibe"] == 10.0
    assert law["top_session"] == ""          # nobody answered q11 in the seed
    assert by_dept["Department of Commerce"]["vibe"] == 6.0

    # Ranked by vibe, best first.
    assert [row["dept"] for row in body["departments"]][0] == "Department of Law"
    assert body["overall"]["filled"] == 2
    assert body["overall"]["vibe"] == 8.0


@pytest.mark.asyncio
async def test_department_analysis_is_scoped_to_the_campus(client):
    await _seed()
    body = (await client.get("/admin/api/orientation/departments",
                             params={"campus": "Kochi"})).json()
    assert [row["dept"] for row in body["departments"]] == ["Department of Law"]
    assert body["departments"][0]["filled"] == 1


@pytest.mark.asyncio
async def test_ppt_download_is_a_real_deck(client):
    await _seed()
    r = await client.get("/admin/survey/orientation-ppt", params={"campus": "Bangalore"})
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    assert "Bangalore" in r.headers["content-disposition"]

    from io import BytesIO
    from pptx import Presentation

    deck = Presentation(BytesIO(r.content))
    assert len(deck.slides) >= 8
    words = " ".join(
        shape.text_frame.text
        for slide in deck.slides for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Student Experience Analysis Report" in words
    assert "Bangalore" in words
    assert "8.0 / 10" in words          # the cohort's average vibe, on the cover
    assert "Loving it".upper() in words.upper()
    # The deck follows the printed report's running order.
    for section in ("Section I", "Section II", "Section III", "Section IV"):
        assert section in words
    assert "RESPONSE OVERVIEW" in words
    assert "Department Scoreboard" in words


@pytest.mark.asyncio
async def test_ppt_download_survives_an_empty_cohort(client):
    await _seed()
    r = await client.get("/admin/survey/orientation-ppt",
                         params={"campus": "Bangalore", "dept": "Department of Design"})
    assert r.status_code == 200
    assert len(r.content) > 10_000


@pytest.mark.asyncio
async def test_ppt_and_department_analysis_need_the_admin_cookie(app_with_mock):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        assert (await anon.get("/admin/api/orientation/departments")).status_code == 403
        assert (await anon.get("/admin/survey/orientation-ppt")).status_code == 403


def test_mood_words_track_the_average():
    from app.orientation_charts import mood_for

    assert mood_for(9.4)[0] == "Buzzing"
    assert mood_for(8.2)[0] == "Loving it"
    assert mood_for(7.0)[0] == "Good vibes"
    assert mood_for(6.4)[0] == "Warming up"
    assert mood_for(5.1)[0] == "Mixed feelings"
    assert mood_for(3.0)[0] == "Needs a lift"
    assert mood_for(None)[0] == "No answers yet"


def test_chart_labels_drop_emoji_and_trim():
    from app.orientation_charts import clean

    assert clean("🎪 Student Club Fair") == "Student Club Fair"
    assert clean("💻 ERP / LMS Onboarding") == "ERP / LMS Onboarding"
    assert clean("x" * 50).endswith("…")


def test_charts_write_pngs_even_with_no_answers(tmp_path):
    from app.orientation_charts import (
        plot_dept_series, plot_nps_ring, plot_response_rate, plot_top_options,
        plot_vibe_hero,
    )

    for name, fn, args in [
        ("vibe", plot_vibe_hero, ({"options": [], "avg": None},)),
        ("nps", plot_nps_ring, ({"promoters": 0, "passives": 0, "detractors": 0},)),
        ("dept", plot_dept_series, ([], [("vibe", "Overall vibe")])),
        ("rate", plot_response_rate, ([],)),
    ]:
        out = fn(*args, tmp_path / f"{name}.png")
        assert out.exists() and out.stat().st_size > 0

    out = plot_top_options([{"label": "🎪 Club fair", "count": 3, "pct": 60.0}],
                           tmp_path / "top.png", "Top sessions")
    assert out.exists() and out.stat().st_size > 0


@pytest.mark.asyncio
async def test_mail_queues_one_message_per_student(client, monkeypatch):
    await _seed()

    sent: list[tuple[str, str]] = []

    class FakeSender:
        async def __aenter__(self): return self
        async def __aexit__(self, *exc): return None
        async def send(self, msg): sent.append((msg["To"], msg["Subject"]))

    from app import emailer
    monkeypatch.setattr(emailer, "SmtpBatchSender", FakeSender)

    r = await client.post(
        "/admin/api/orientation/mail",
        params={"campus": "Bangalore", "group": "filled"},
        json={"subject": "Thank you", "message": "First para.\n\nSecond para."},
    )
    assert r.status_code == 200
    assert r.json()["total"] == 2

    # BackgroundTasks run once the response is delivered.
    assert [subject for _, subject in sent] == ["Thank you", "Thank you"]
    assert {addr for addr, _ in sent} == {"Asha <a@x.com>", "Bilal <b@x.com>"}

    task = await db.get_db()["admin_tasks"].find_one({"_id": r.json()["task_id"]})
    assert (task["status"], task["sent"], task["failed"]) == ("completed", 2, 0)

    user = await db.get_db()["users"].find_one({"email": "a@x.com"})
    assert user["orientation_mail_count"] == 1


@pytest.mark.asyncio
async def test_mail_needs_a_subject_and_a_body(client):
    await _seed()
    r = await client.post("/admin/api/orientation/mail",
                          params={"campus": "Bangalore"}, json={"subject": "Hi"})
    assert r.status_code == 400


@pytest.mark.asyncio
async def test_mail_refuses_an_empty_cohort(client):
    await _seed()
    r = await client.post("/admin/api/orientation/mail",
                          params={"campus": "Kochi", "group": "pending"},
                          json={"subject": "Hi", "message": "There"})
    assert r.status_code == 400
    assert r.json()["total"] == 0


@pytest.mark.asyncio
async def test_orientation_endpoints_need_the_survey_admin_cookie(app_with_mock):
    transport = ASGITransport(app=app_with_mock)
    async with AsyncClient(transport=transport, base_url="http://test") as anon:
        for path in ("/admin/api/orientation/campuses",
                     "/admin/api/orientation/report",
                     "/admin/api/orientation/students"):
            assert (await anon.get(path)).status_code == 403
        assert (await anon.post("/admin/api/orientation/mail",
                                json={"subject": "a", "message": "b"})).status_code == 403

    # The orientation admin's cookie is not enough either — this report lives in
    # the survey admin.
    async with AsyncClient(transport=transport, base_url="http://test") as ori:
        ori.cookies.set("orientation_admin_session", "1")
        assert (await ori.get("/admin/api/orientation/report")).status_code == 403


@pytest.mark.asyncio
async def test_mail_body_becomes_paragraphs_and_carries_a_resume_link():
    from app import emailer

    msg = emailer.build_orientation_message(
        "a@x.com", "Asha", "Subject here", "First para.\n\nSecond para.",
        link="http://test/resume/abc?src=reminder", link_label="Open my survey",
        campus="Bangalore",
    )
    html = ""
    text = ""
    for part in msg.walk():
        if part.get_content_type() == "text/html":
            html = part.get_payload(decode=True).decode()
        if part.get_content_type() == "text/plain":
            text = part.get_payload(decode=True).decode()

    assert msg["To"] == "Asha <a@x.com>"
    assert "<p" in html and "First para." in html and "Second para." in html
    assert "http://test/resume/abc?src=reminder" in html
    assert "Bangalore" in html
    assert text.startswith("Dear Asha,")
    assert "Open my survey: http://test/resume/abc?src=reminder" in text


# ── What the deck is built from ──────────────────────────────────────────────
@pytest.mark.asyncio
async def test_campus_split_counts_ug_and_pg_per_campus(app_with_mock):
    """The deck's opening slide: who answered, campus by campus."""
    await _seed()
    from app.orientation_data import campus_split, orientation_dataset

    data = await orientation_dataset()
    rows = {r["campus"]: r for r in campus_split(data["filled"], data["pending"])}

    assert rows["Bangalore"]["filled"] == 2          # Asha and Bilal
    assert rows["Bangalore"]["pending"] == 1         # Dev owes the form
    assert rows["Bangalore"]["ug"] == 2 and rows["Bangalore"]["pg"] == 0
    assert rows["Bangalore"]["departments"] == 2
    assert rows["Kochi"]["filled"] == 1
    # Esha never finished the baseline, so she is in neither column.
    assert sum(r["eligible"] for r in rows.values()) == 4


@pytest.mark.asyncio
async def test_department_rows_carry_the_nps_denominator(app_with_mock):
    """Promoters mean nothing without the count they are a share of.

    A department where one student answered warmly and nobody coldly is not a
    department where 100% are promoters — it is one where one student answered.
    """
    await _seed()
    from app.orientation_data import department_rows, orientation_dataset

    data = await orientation_dataset(campus="Bangalore")
    rows = {r["dept"]: r for r in department_rows(data["filled"], data["pending"])}

    law = rows["Department of Law"]
    assert law["promoters"] == 1 and law["detractors"] == 0
    assert law["nps_answered"] == 1          # Asha; Dev never answered
    commerce = rows["Department of Commerce"]
    assert commerce["detractors"] == 1 and commerce["nps_answered"] == 1

