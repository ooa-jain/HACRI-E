"""
orientation_ppt.py — the Deeksharambh orientation report as a slide deck.

Built to the house design of the Deeksharambh analysis decks: a mint ground
with a white panel floating on it, navy serif headings centred and underlined,
a section kicker above each one ("Section I — PROGRAM EFFECTIVENESS"), teal
meters for the averages, and a department-wise chart with the observations
listed beside it.

Everything comes from `summarize_orientation()`, so the deck and the dashboard
can never tell different stories. Where the deck says something in words — the
observations beside a chart — it is composed from those same figures and never
from anything else.
"""
from __future__ import annotations

import io
import tempfile
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.orientation_analysis import MIN_REPORTABLE
from app.orientation_charts import (
    clean, mood_for, plot_dept_series, plot_nps_ring, plot_response_rate,
    plot_top_options, plot_vibe_hero,
)

# ── The palette of the printed report ────────────────────────────────────────
NAVY = RGBColor(0x2E, 0x3A, 0x64)     # every heading
INK = RGBColor(0x2F, 0x36, 0x40)      # body copy
MUTED = RGBColor(0x5B, 0x65, 0x70)
TEAL = RGBColor(0x21, 0xA8, 0x8A)     # meters, figures, rules
TEAL_DEEP = RGBColor(0x17, 0x80, 0x6A)
MINT = RGBColor(0xD9, 0xF1, 0xEC)     # the ground
MINT_DEEP = RGBColor(0xBF, 0xE7, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xDF, 0xE4, 0xE8)
GOLD = RGBColor(0xC8, 0xA4, 0x5A)

SERIF = "Georgia"

W = Inches(13.333)
H = Inches(7.5)
RECT = 1  # MSO_SHAPE.RECTANGLE
ROUND = 5  # MSO_SHAPE.ROUNDED_RECTANGLE

# Scoreboard rows per slide — what fits above the footnote at this type size.
PER_PAGE = 12


def _hex(colour: str) -> RGBColor:
    return RGBColor.from_string(colour.lstrip("#").upper())


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _face(para, *, size, bold, colour, italic=False, underline=False) -> None:
    """Set the type on a paragraph and on the runs inside it.

    PowerPoint reads a run's own properties first and only falls back to the
    paragraph's, so a deck that sets one and not the other renders in whatever
    the theme fancies.
    """
    for font in [para.font] + [run.font for run in para.runs]:
        font.name = SERIF
        font.size = Pt(size)
        font.bold = bold
        font.italic = italic
        font.underline = underline
        font.color.rgb = colour


def _text(slide, left, top, width, height, text, *, size=14, bold=False,
          colour=INK, align=PP_ALIGN.LEFT, italic=False, underline=False,
          spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.text = text
    para.alignment = align
    if spacing:
        para.line_spacing = spacing
    _face(para, size=size, bold=bold, colour=colour, italic=italic, underline=underline)
    return box


def _paras(slide, left, top, width, height, lines: list[tuple], *, spacing=1.25):
    """A block of paragraphs: (text, size, bold, colour) each, or a blank line."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        text, size, bold, colour = line
        para.text = text
        para.line_spacing = spacing
        _face(para, size=size, bold=bold, colour=colour)
    return box


def _rect(slide, left, top, width, height, fill: RGBColor | None,
          line: RGBColor | None = None, shape=RECT, line_width=Pt(1)):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = line_width
    box.shadow.inherit = False
    return box


def _stage(slide, *, panel=True, border: RGBColor | None = None):
    """The mint ground with the white panel floating on it.

    Every slide in the printed report is built this way; the panel is dropped
    only where a chart wants the full mint width.
    """
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = MINT
    _rect(slide, Emu(0), Emu(0), W, Inches(1.05), WHITE)
    _rect(slide, Emu(0), H - Inches(0.62), W, Inches(0.62), WHITE)
    if panel:
        _rect(slide, Inches(0.28), Inches(1.28), W - Inches(0.56),
              H - Inches(2.15), WHITE, border)


def _title(slide, title: str, kicker: str = "", *, top=Inches(1.5),
           underline=True, size=32):
    """Centred kicker over a centred serif heading, ruled underneath."""
    if kicker:
        _text(slide, Inches(0.8), top, W - Inches(1.6), Inches(0.34),
              kicker, size=13, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)
        top = top + Inches(0.4)
    _text(slide, Inches(0.8), top, W - Inches(1.6), Inches(0.62),
          title, size=size, bold=True, colour=NAVY, align=PP_ALIGN.CENTER,
          underline=underline)
    return top + Inches(0.75)


def _meter(slide, left, top, width, *, value: float | None, maximum: int,
           label: str, note: str):
    """A rounded teal bar with its score beside it, then the label and a note."""
    bar_w = width - Inches(1.55)
    track = _rect(slide, left, top, bar_w, Inches(0.22), MINT_DEEP, shape=ROUND)
    track.adjustments[0] = 0.5
    share = 0.0 if value is None else max(0.0, min(1.0, value / maximum))
    if share > 0:
        fill = _rect(slide, left, top, Emu(int(bar_w * max(share, 0.06))),
                     Inches(0.22), TEAL, shape=ROUND)
        fill.adjustments[0] = 0.5
    _text(slide, left + bar_w + Inches(0.16), top - Inches(0.09), Inches(1.35), Inches(0.3),
          "—" if value is None else f"{value:.2f}/{maximum}",
          size=15, bold=True, colour=INK)
    _text(slide, left, top + Inches(0.28), width, Inches(0.3),
          label, size=15, bold=True, colour=NAVY)
    _text(slide, left, top + Inches(0.62), width, Inches(0.3),
          note, size=12, colour=MUTED)


def _tile(slide, left, top, width, height, title: str, body: str):
    """A mint card with a teal dot — the engagement slide is built from these."""
    _rect(slide, left, top, width, height, MINT, MINT_DEEP, shape=ROUND)
    dot = _rect(slide, left + Inches(0.22), top + Inches(0.2),
                Inches(0.42), Inches(0.42), TEAL, shape=ROUND)
    dot.adjustments[0] = 0.5
    _text(slide, left + Inches(0.22), top + Inches(0.78), width - Inches(0.44),
          Inches(0.34), title, size=14, bold=True, colour=NAVY)
    _text(slide, left + Inches(0.22), top + Inches(1.12), width - Inches(0.44),
          Inches(0.5), body, size=12.5, colour=INK)


def _ranked(slide, left, top, width, height, title: str, note: str,
            options: list[dict], answered: int, *, accent=TEAL_DEEP, limit=5):
    """One feedback question as a ranked list, with the number behind it.

    The count is what students actually clicked; the share beside it is of the
    students who answered that question, which the note above the list names
    outright. A ranked list with no denominator is the one way this deck could
    still mislead a reader who is doing nothing wrong.
    """
    _rect(slide, left, top, width, height, MINT, MINT_DEEP, shape=ROUND)
    _text(slide, left + Inches(0.2), top + Inches(0.16), width - Inches(0.4),
          Inches(0.3), title, size=14, bold=True, colour=accent)
    _text(slide, left + Inches(0.2), top + Inches(0.46), width - Inches(0.4),
          Inches(0.26), note, size=10.5, colour=MUTED)

    rows = (options or [])[:limit]
    if not rows:
        _text(slide, left + Inches(0.2), top + Inches(0.82), width - Inches(0.4),
              Inches(0.3), "Nobody answered this question.", size=11.5, colour=MUTED)
        return
    for i, option in enumerate(rows):
        y = top + Inches(0.8) + i * Inches(0.46)
        _text(slide, left + Inches(0.2), y, Inches(0.3), Inches(0.3),
              f"{i + 1}", size=11.5, bold=True, colour=accent)
        _text(slide, left + Inches(0.5), y, width - Inches(1.85), Inches(0.3),
              clean(option.get("label", ""), 38), size=11.5, colour=INK)
        _text(slide, left + width - Inches(1.35), y, Inches(1.15), Inches(0.3),
              f"{option.get('count', 0):,}  ·  {option.get('pct', 0):.0f}%",
              size=11.5, bold=True, colour=NAVY, align=PP_ALIGN.RIGHT)


def _observations(slide, left, top, width, height, bullets: list[str],
                  *, fill=MINT, border=MINT_DEEP):
    """The commentary box beside a department chart."""
    bullets = [b for b in bullets if b]
    if not bullets:
        return
    _rect(slide, left, top, width, height, fill, border)
    box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.16),
                                   width - Inches(0.32), height - Inches(0.32))
    frame = box.text_frame
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = f"•  {bullet}"
        para.line_spacing = 1.2
        para.space_after = Pt(9)
        _face(para, size=11.5, bold=False, colour=INK)


def _table(slide, left, top, width, headers: list[str], rows: list[list[str]],
           widths: list[float] | None = None):
    """A quiet table: mint header row, hairline rules, figures on the right."""
    if not rows:
        return
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width,
                                   Inches(0.26) * (len(rows) + 1))
    table = shape.table
    table.first_row = True
    table.horz_banding = False
    if widths:
        for i, share in enumerate(widths):
            table.columns[i].width = Emu(int(width * share))
    for col, label in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = MINT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if col else PP_ALIGN.LEFT
        _face(para, size=10.5, bold=True, colour=NAVY)
    for r, values in enumerate(rows, start=1):
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            _face(para, size=10, bold=False, colour=INK)


def _picture(slide, image: Path, left, top, max_w, max_h) -> None:
    """Drop a chart in, scaled to fit its box and centred in it.

    Charts are as tall as their data needs (one group per department), so a
    fixed width alone would push the long ones off the bottom of the slide.
    """
    if not image or not Path(image).exists():
        return
    from PIL import Image

    with Image.open(image) as img:
        px_w, px_h = img.size
    scale = min(max_w / px_w, max_h / px_h)
    width, height = int(px_w * scale), int(px_h * scale)
    slide.shapes.add_picture(
        str(image),
        left + int((max_w - width) / 2),
        top + int((max_h - height) / 2),
        width=width, height=height,
    )


def _fmt(value, digits=1, suffix="") -> str:
    return "—" if value is None else f"{value:.{digits}f}{suffix}"


def _top(options, index: int = 0) -> dict:
    rows = options or []
    return rows[index] if len(rows) > index else {}


def _leading(options) -> str:
    """The answer most students gave, with its share — "Yes, mostly (64%)".

    Deliberately quotes the option rather than bucketing it: the form's wording
    changes between years, and a matcher that decides "Somewhat easy" counts as
    easy would overstate the result in the one direction nobody would check.
    """
    rows = [o for o in (options or []) if o.get("count")]
    if not rows:
        return ""
    top = rows[0]
    return f"{clean(top['label'], 40)} ({top.get('pct', 0):.0f}%)"


def _spread(rows: list[dict], key: str) -> tuple[dict, dict] | None:
    """The highest and lowest department on one measure, when there are two.

    Only departments big enough to report: naming "the department that rated
    the week highest" is worth saying about 91 students and meaningless about
    one, who would top the table on their own mood that morning.
    """
    scored = [r for r in rows
              if r.get(key) is not None and r.get("reportable", True)]
    if len(scored) < 2:
        return None
    ranked = sorted(scored, key=lambda r: -r[key])
    return ranked[0], ranked[-1]


def _of(answered, total) -> str:
    """"of the 1,078 who answered" — the denominator a bare share hides.

    Every percentage in this deck is a share of the students who answered that
    one question, which is almost never the whole cohort. Printed next to the
    figure so nobody reads it against the response count on the cover.
    """
    if not answered:
        return ""
    if total and answered >= total:
        return f"of all {answered:,} respondents"
    return f"of the {answered:,} who answered"


def _denominator(answered, total) -> str:
    """The note above a ranked list, naming what its shares are shares of."""
    if not answered:
        return "Nobody answered this question."
    if total and answered >= total:
        return f"All {answered:,} respondents answered"
    return (f"{answered:,} of {total:,} respondents answered · "
            f"shares are of those {answered:,}")


def _rule(slide, left, top, width) -> None:
    _rect(slide, left, top, width, Emu(9525), LINE, None)


def generate_orientation_ppt(
    *,
    campus: str,
    scope: str,
    report: dict,
    departments: list[dict],
    campuses: list[dict] | None = None,
) -> bytes:
    """Build the deck. `report` is a build_report() / summarize_orientation() result."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    head = report.get("headline", {})
    cover = report.get("coverage", {})
    highlights = report.get("highlights", {})
    levels = report.get("levels", [])
    dept_mix = report.get("departments", [])
    campuses = campuses or []
    questions = {q["key"]: q for s in report.get("sections", []) for q in s["questions"]}
    total = report.get("count", 0)
    mood_word, mood_hex = mood_for(head.get("vibe"))

    ug = next((l["count"] for l in levels if str(l.get("level", "")).upper() == "UG"), 0)
    pg = next((l["count"] for l in levels if str(l.get("level", "")).upper() == "PG"), 0)

    # ── Slide 1 · Cover ───────────────────────────────────────────────────────
    slide = _blank(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    _rect(slide, Emu(0), Inches(1.15), Inches(8.0), Inches(5.35), MINT)
    _rect(slide, Inches(8.0), Inches(1.15), W - Inches(8.0), Inches(5.35), MINT_DEEP)

    _text(slide, Inches(0.62), Inches(1.85), Inches(6.9), Inches(1.9),
          f"Student Experience Analysis Report — Deeksharambh {datetime.now():%Y}",
          size=34, bold=True, colour=NAVY, spacing=1.15)
    _text(slide, Inches(0.62), Inches(4.0), Inches(6.6), Inches(1.2),
          "This report analyses how the first week landed — the transition into "
          "campus life, the Bridge Course, the orientation sessions themselves, "
          "and what students say they now expect of us.",
          size=12.5, colour=INK, spacing=1.35)
    _text(slide, Inches(0.62), Inches(5.35), Inches(6.6), Inches(0.4),
          f"{campus}  ·  {scope}", size=12.5, bold=True, colour=NAVY)
    _text(slide, Inches(0.62), Inches(5.72), Inches(6.6), Inches(0.4),
          f"Office of Academics  ·  {total} student responses  ·  "
          f"{datetime.now():%d %B %Y}", size=11.5, colour=MUTED)

    # The right panel carries the one number the room came for.
    _text(slide, Inches(8.5), Inches(2.25), Inches(4.4), Inches(0.4),
          "OVERALL VIBE OF THE STUDENTS", size=11.5, bold=True, colour=TEAL_DEEP,
          align=PP_ALIGN.CENTER)
    _text(slide, Inches(8.5), Inches(2.7), Inches(4.4), Inches(1.1),
          _fmt(head.get("vibe"), 1, " / 10"), size=54, bold=True, colour=NAVY,
          align=PP_ALIGN.CENTER)
    _text(slide, Inches(8.5), Inches(3.85), Inches(4.4), Inches(0.4),
          mood_word.upper(), size=15, bold=True, colour=_hex(mood_hex),
          align=PP_ALIGN.CENTER)
    for i, (value, label) in enumerate([
        (_fmt(head.get("nps"), 0), "Net Promoter Score"),
        (_fmt(head.get("belonging"), 1, " / 10"), "Sense of belonging"),
        (_fmt(cover.get("pct"), 0, "%"), "Response rate"),
    ]):
        top = Inches(4.45) + i * Inches(0.62)
        _text(slide, Inches(8.5), top, Inches(2.6), Inches(0.34),
              label, size=12, colour=INK)
        _text(slide, Inches(11.1), top, Inches(1.8), Inches(0.34),
              value, size=12, bold=True, colour=NAVY, align=PP_ALIGN.RIGHT)

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)

        # ── Slide 2 · Response overview ───────────────────────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, "RESPONSE OVERVIEW", top=Inches(1.6), size=30)

        campus_lines: list[tuple] = [("Campus-wise Split:", 15, True, NAVY)]
        for row in campuses or [{"campus": campus, "filled": total,
                                 "ug": ug, "pg": pg,
                                 "departments": len(dept_mix)}]:
            campus_lines.append((
                f"•  {row['campus']}: {row['filled']} responses "
                f"(UG: {row.get('ug', 0)}, PG: {row.get('pg', 0)})", 13, False, INK))
            campus_lines.append((
                f"     from {row.get('departments', 0)} departments", 12, False, MUTED))

        _text(slide, Inches(0.75), Inches(2.72), Inches(5.6), Inches(0.4),
              f"Total Responses: {total}", size=16, bold=True, colour=TEAL_DEEP)
        _paras(slide, Inches(0.75), Inches(3.2), Inches(5.6), Inches(1.6), campus_lines)
        _text(slide, Inches(0.75), Inches(4.75), Inches(5.6), Inches(0.34),
              "Departments with the highest response", size=13, bold=True, colour=NAVY)
        _table(slide, Inches(0.75), Inches(5.12), Inches(5.6),
               ["Department", "Responses", "Share"],
               [[clean(r["dept"], 40), str(r["count"]), f"{r.get('pct', 0):.0f}%"]
                for r in dept_mix[:6]],
               widths=[0.58, 0.22, 0.20])

        _text(slide, Inches(6.95), Inches(2.72), Inches(5.6), Inches(0.4),
              "Overall UG vs PG Breakdown:", size=16, bold=True, colour=NAVY)
        _paras(slide, Inches(6.95), Inches(3.2), Inches(5.6), Inches(0.9), [
            (f"•  UG Students: {ug} responses", 13, False, INK),
            (f"•  PG Students: {pg} responses", 13, False, INK),
        ])
        # `dept_mix` also carries the bucket for replies we could not match to a
        # student record, which is not a department. Counting it here is what
        # used to make this line disagree with the campus split above it.
        named = report.get("department_count", len(dept_mix))
        unmatched = report.get("unmatched", 0)
        _text(slide, Inches(6.95), Inches(4.02), Inches(5.6), Inches(0.7),
              f"{cover.get('filled', 0)} of {cover.get('eligible', 0)} students in scope "
              f"have answered ({_fmt(cover.get('pct'), 0, '%')}), across "
              f"{named} departments."
              + (f" A further {unmatched} replies could not be matched to a "
                 "student record and are counted in the totals only."
                 if unmatched else ""),
              size=12.5, colour=INK, spacing=1.35)
        _text(slide, Inches(6.95), Inches(4.75), Inches(5.6), Inches(0.34),
              "Departments still to be heard from", size=13, bold=True, colour=NAVY)
        _table(slide, Inches(6.95), Inches(5.12), Inches(5.6),
               ["Department", "Answered", "Pending"],
               [[clean(r["dept"], 40), str(r.get("filled", 0)), str(r.get("pending", 0))]
                for r in sorted(departments, key=lambda r: -r.get("pending", 0))[:6]
                if r.get("pending", 0)],
               widths=[0.58, 0.21, 0.21])

        # ── Slides 3-4 · What students told us ────────────────────────────────
        #
        # First, before any average. Every one of these questions was asked so
        # the programme could be changed, and a deck that opens on scores and
        # leaves the asks to an appendix has quietly reordered the point. The
        # dashboard has always drawn all six panels; the deck used to print
        # four and drop "stop" and "introduce" entirely.
        answered_by = report.get("highlights_answered", {})
        labels = report.get("highlights_labels", {})

        def feedback_slide(heading, kicker, panels):
            slide = _blank(prs)
            _stage(slide)
            _title(slide, heading, kicker, top=Inches(1.4), size=26)
            for i, (key, accent) in enumerate(panels):
                answered = answered_by.get(key, 0)
                _ranked(slide,
                        Inches(0.62) + i * Inches(4.08), Inches(2.72),
                        Inches(3.85), Inches(3.35),
                        labels.get(key, key),
                        _denominator(answered, total),
                        highlights.get(key) or [], answered, accent=accent)
            return slide

        slide = feedback_slide(
            "What students asked us to change",
            "Section 0 — STUDENT FEEDBACK, IN FULL",
            [("keep", TEAL_DEEP), ("stop", _hex("#b03a5b")), ("introduce", _hex("#6f6bd8"))])
        _text(slide, Inches(0.62), Inches(6.22), Inches(12.1), Inches(0.5),
              ("Every question here was a fixed list of options — the orientation form "
               "carries no free-text box, so these rankings are the whole of what "
               "students were able to tell us in their own direction.  "
               f"Answered by {answered_by.get('keep', 0):,}, "
               f"{answered_by.get('stop', 0):,} and "
               f"{answered_by.get('introduce', 0):,} students respectively, "
               f"of {total:,} who responded."),
              size=11.5, colour=MUTED, spacing=1.25)

        slide = feedback_slide(
            "Where the week was hardest",
            "Section 0 — STUDENT FEEDBACK, IN FULL",
            [("challenges", _hex("#c0504d")), ("stressors", _hex("#e0913a")),
             ("least_connecting", _hex("#5b6b8c"))])
        reasons = highlights.get("reasons") or []
        _text(slide, Inches(0.62), Inches(6.22), Inches(12.1), Inches(0.5),
              (f"{labels.get('reasons', 'Why they scored us that way')}: "
               + "; ".join(f"{clean(o['label'], 38)} ({o['count']:,})"
                           for o in reasons[:3])
               + f" — {_of(answered_by.get('reasons', 0), total)}."
               if reasons else
               f"{labels.get('reasons', 'Why they scored us that way')}: "
               "nobody answered that question."),
              size=11.5, colour=MUTED, spacing=1.25)

        # ── Slide 5 · Participation, department by department ─────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, f"Department wise — Who has answered — {campus}",
               top=Inches(1.5), size=23)
        rate_img = plot_response_rate(departments, tmpdir / "rate.png", title="")
        _picture(slide, rate_img, Inches(0.55), Inches(2.35), Inches(8.4), Inches(4.1))
        best_rate = max((r for r in departments if r.get("eligible")),
                        key=lambda r: r.get("pct", 0), default=None)
        _observations(slide, Inches(9.2), Inches(2.5), Inches(3.5), Inches(3.1), [
            f"{cover.get('filled', 0)} of {cover.get('eligible', 0)} eligible students "
            f"answered — {_fmt(cover.get('pct'), 0, '%')} of the cohort.",
            (f"{clean(best_rate['dept'], 44)} answered most completely "
             f"({best_rate.get('pct', 0):.0f}%)." if best_rate else ""),
            (f"{cover.get('pending', 0)} students have finished the baseline but not "
             "yet the orientation form." if cover.get("pending") else ""),
        ])

        # ── Slide 4 · Section I — program effectiveness ────────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, "Ease of Transition to Campus Life",
               "Section I — PROGRAM EFFECTIVENESS", top=Inches(1.55))
        _rect(slide, Inches(0.62), Inches(2.82), W - Inches(1.24), Inches(2.98),
              None, TEAL)

        welcomed = _leading(questions.get("q3", {}).get("options"))
        easy = _leading(questions.get("q5", {}).get("options"))
        knows_help = _leading(questions.get("q31", {}).get("options"))
        vibe_q = questions.get("q2", {})
        high_vibe = sum(o.get("pct", 0) for o in vibe_q.get("options", [])
                        if o.get("count") and int(o["label"]) >= 8)

        for i, (value, maximum, label, note) in enumerate([
            (head.get("vibe"), 10, "Overall vibe of the week",
             f"{high_vibe:.0f}% of the {vibe_q.get('answered', 0):,} students who "
             f"rated the week gave it 8 or more."
             if vibe_q.get("options") else "Not enough ratings yet."),
            (head.get("belonging"), 10, "Sense of belonging at JAIN",
             "How much students already feel part of the place."),
            (head.get("success"), 10, "Confidence of succeeding here",
             "Whether they can see themselves doing well."),
            (head.get("nps_avg"), 10, "Likelihood to recommend JAIN",
             f"Net Promoter Score of {_fmt(head.get('nps'), 0)} across "
             f"{head.get('nps_answered', 0)} answers."),
        ]):
            left = Inches(0.85) + (i % 2) * Inches(6.1)
            top = Inches(3.05) + (i // 2) * Inches(1.35)
            _meter(slide, left, top, Inches(5.5), value=value, maximum=maximum,
                   label=label, note=note)

        # Each clause quotes the answer most students actually chose, so the
        # caption cannot round a mixed result up into a good one.
        # Each of these three questions was answered by a different number of
        # students, so each share is of its own question's respondents. Said
        # once, plainly, rather than three parenthetical denominators.
        summary = []
        if welcomed:
            summary.append(f"felt welcomed — {welcomed}")
        if easy:
            summary.append(f"transition into university life — {easy}")
        if knows_help:
            summary.append(f"knows whom to ask for help — {knows_help}")
        _text(slide, Inches(0.85), Inches(5.93), Inches(11.6), Inches(0.6),
              ("The first week averaged " + _fmt(head.get("vibe"), 1) +
               " out of 10 — " + mood_word.lower() + "."
               + ("  Most common answers: " + "; ".join(summary)
                  + ".  Each share is of the students who answered that "
                    "question, not of all " + f"{total:,}."
                  if summary else "")),
              size=12.5, colour=INK, spacing=1.25)

        # ── Slide 5 · Section I, department by department ──────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, f"Department wise — Ease of Transition to Campus Life — {campus}",
               top=Inches(1.5), size=22)
        series_img = plot_dept_series(
            departments,
            [("vibe", "Overall vibe"), ("belonging", "Sense of belonging"),
             ("success", "Confidence of succeeding")],
            tmpdir / "dept_transition.png", maximum=10)
        _picture(slide, series_img, Inches(0.5), Inches(2.35), Inches(8.5), Inches(4.1))
        spread = _spread(departments, "vibe")
        _observations(slide, Inches(9.2), Inches(2.5), Inches(3.5), Inches(3.3), [
            (f"{clean(spread[0]['dept'], 42)} rated the week highest at "
             f"{spread[0]['vibe']:.1f}/10; {clean(spread[1]['dept'], 42)} lowest at "
             f"{spread[1]['vibe']:.1f}/10." if spread else
             "Only one department has answered so far."),
            f"The cohort averages {_fmt(head.get('vibe'), 1)}/10 on vibe and "
            f"{_fmt(head.get('belonging'), 1)}/10 on belonging.",
            (f"Departments answering most: " +
             ", ".join(f"{clean(r['dept'], 30)} ({r['count']})" for r in dept_mix[:3])
             if dept_mix else ""),
        ])

        # ── Slide 6 · Section II — academic foundation ─────────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, "Bridge Course: Foundational Learning Modules",
               "Section II — ACADEMIC FOUNDATION", top=Inches(1.55), size=28)

        prepared = _leading(questions.get("q18", {}).get("options"))
        bridge_n = questions.get("q16", {}).get("answered", 0)
        prepared_n = questions.get("q18", {}).get("answered", 0)
        _paras(slide, Inches(0.85), Inches(3.0), Inches(6.4), Inches(2.0), [
            (f"•  Academic confidence after the Bridge Course: "
             f"{_fmt(head.get('bridge'), 2)}/5", 14, True, NAVY),
            (f"•  Rated by {bridge_n:,} of the {total:,} who responded",
             12.5, False, MUTED),
            (f"•  Prepared for regular classes: {prepared or '—'}", 14, True, NAVY),
            (f"•  That share is of the {prepared_n:,} who answered that question",
             12.5, False, MUTED),
        ])
        helpful = (questions.get("q17", {}).get("options") or [])[:4]
        _text(slide, Inches(0.85), Inches(4.35), Inches(6.4), Inches(1.4),
              ("The Bridge Course was built to get students ready for the way the "
               "programme actually teaches. These are the figures behind that, and "
               "the areas students named as the ones that helped."),
              size=12.5, colour=INK, spacing=1.35)

        bridge_img = plot_top_options(
            questions.get("q17", {}).get("options") or [],
            tmpdir / "helpful.png", "Bridge Course areas that helped most", limit=5)
        _picture(slide, bridge_img, Inches(7.15), Inches(2.85), Inches(5.4), Inches(3.6))

        # ── Slide 7 · Section II, department by department ─────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide,
               f"Department wise — Bridge Course confidence — {campus}",
               top=Inches(1.5), size=22)
        bridge_dept = plot_dept_series(
            departments, [("bridge", "Academic confidence after the Bridge Course")],
            tmpdir / "dept_bridge.png", maximum=5,
            empty="No Bridge Course ratings yet")
        _picture(slide, bridge_dept, Inches(0.5), Inches(2.35), Inches(8.5), Inches(4.1))
        bridge_spread = _spread(departments, "bridge")
        _observations(slide, Inches(9.2), Inches(2.5), Inches(3.5), Inches(3.1), [
            f"The Bridge Course averages {_fmt(head.get('bridge'), 2)} out of 5 "
            "across the cohort.",
            (f"{clean(bridge_spread[0]['dept'], 42)} rated it highest "
             f"({bridge_spread[0]['bridge']:.1f}/5); "
             f"{clean(bridge_spread[1]['dept'], 42)} lowest "
             f"({bridge_spread[1]['bridge']:.1f}/5)." if bridge_spread else ""),
            (f"Most helpful: {clean(_top(helpful).get('label', ''), 44)}."
             if helpful else ""),
        ])

        # ── Slide 8 · Section III — engagement and networking ──────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, "Impact of Immersive Orientation",
               "Section III — ENGAGEMENT AND NETWORKING", top=Inches(1.5))
        top_sessions = (highlights.get("impactful") or [])[:3]
        impact_n = answered_by.get("impactful", 0)
        if top_sessions:
            for i, option in enumerate(top_sessions):
                _tile(slide, Inches(0.85) + i * Inches(4.0), Inches(3.0),
                      Inches(3.7), Inches(1.85),
                      clean(option["label"], 34),
                      f"{option['count']:,} students named it among the sessions "
                      f"with the biggest impact — {option.get('pct', 0):.0f}% "
                      f"{_of(impact_n, total)}.")
        else:
            _text(slide, Inches(0.85), Inches(3.2), Inches(11.6), Inches(0.5),
                  "No sessions have been named yet.", size=13, colour=MUTED)

        _text(slide, Inches(0.85), Inches(5.15), Inches(11.6), Inches(1.1),
              (f"{head.get('promoters', 0)} of {head.get('nps_answered', 0)} students "
               f"who answered would actively recommend JAIN, against "
               f"{head.get('detractors', 0)} who would not — a Net Promoter Score of "
               f"{_fmt(head.get('nps'), 0)}. "
               + (f"The session most often named for improvement was "
                  f"{clean(_top(highlights.get('needs_work')).get('label', ''), 44)}."
                  if highlights.get("needs_work") else "")),
              size=13, colour=INK, spacing=1.3)

        # ── Slide 9 · Section III, department by department ────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, f"Department wise — Would they recommend JAIN — {campus}",
               top=Inches(1.5), size=22)
        for row in departments:
            # Out of everyone in that department who answered the question —
            # promoters against detractors alone would read 100% for a
            # department where one student answered warmly and nobody coldly.
            answered = row.get("nps_answered") or 0
            row["_promoters_pct"] = (100.0 * row.get("promoters", 0) / answered
                                     if answered else None)
            row["_detractors_pct"] = (100.0 * row.get("detractors", 0) / answered
                                      if answered else None)
        nps_dept = plot_dept_series(
            departments,
            [("_promoters_pct", "Promoters (9–10)"), ("_detractors_pct", "Detractors (0–6)")],
            tmpdir / "dept_nps.png", maximum=100,
            empty="No recommendation scores yet")
        _picture(slide, nps_dept, Inches(0.5), Inches(2.35), Inches(8.5), Inches(4.1))
        _observations(slide, Inches(9.2), Inches(2.5), Inches(3.5), Inches(3.1), [
            f"Across the cohort: {head.get('promoters', 0)} promoters, "
            f"{head.get('passives', 0)} passives and {head.get('detractors', 0)} "
            f"detractors — NPS {_fmt(head.get('nps'), 0)}.",
            "Bars are each department's own answers, so a small department's "
            "percentage rests on few students.",
        ])

        # ── Slide 10 · What to keep, what to fix ───────────────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, "What to keep, and what to fix", top=Inches(1.5), size=24)
        keep_img = plot_top_options(highlights.get("keep") or [],
                                    tmpdir / "keep.png", "Keep next year", limit=6)
        fix_img = plot_top_options(highlights.get("needs_work") or [],
                                   tmpdir / "fix.png", "Sessions needing work",
                                   colour="#c0504d", limit=6)
        _picture(slide, keep_img, Inches(0.6), Inches(2.35), Inches(5.9), Inches(4.1))
        _picture(slide, fix_img, Inches(6.85), Inches(2.35), Inches(5.9), Inches(4.1))

        # ── Slide 11 · The vibe, score by score ────────────────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, f"How the week was rated, score by score — {campus}",
               top=Inches(1.5), size=24)
        vibe_img = plot_vibe_hero(
            questions.get("q2", {"options": [], "avg": head.get("vibe")}),
            tmpdir / "vibe.png", title="")
        _picture(slide, vibe_img, Inches(0.6), Inches(2.35), Inches(7.4), Inches(4.1))
        ring = plot_nps_ring(questions.get("q34") or {
            "promoters": head.get("promoters", 0),
            "passives": head.get("passives", 0),
            "detractors": head.get("detractors", 0),
            "nps": head.get("nps"),
        }, tmpdir / "nps.png")
        _picture(slide, ring, Inches(8.2), Inches(2.35), Inches(4.6), Inches(4.1))

        # ── Slide 12 · Section IV — aspirations and growth ─────────────────────
        slide = _blank(prs)
        _stage(slide, border=TEAL)
        _title(slide, "Reflection and Future Readiness",
               "Section IV — ASPIRATIONS AND GROWTH", top=Inches(1.55))

        _text(slide, Inches(0.85), Inches(3.0), Inches(5.5), Inches(0.34),
              "Most helpful aspects:", size=14, bold=True, colour=TEAL_DEEP)
        _paras(slide, Inches(0.85), Inches(3.4), Inches(5.5), Inches(1.5),
               [(f"•  {clean(o['label'], 44)}", 12.5, False, INK)
                for o in (questions.get("q19", {}).get("options") or [])[:4]] or
               [("Nothing named yet.", 12.5, False, MUTED)])

        _text(slide, Inches(6.9), Inches(3.0), Inches(5.5), Inches(0.34),
              "Top student expectations:", size=14, bold=True, colour=TEAL_DEEP)
        _paras(slide, Inches(6.9), Inches(3.4), Inches(5.5), Inches(1.5),
               [(f"•  {clean(o['label'], 44)}", 12.5, False, INK)
                for o in (questions.get("q33", {}).get("options") or [])[:4]] or
               [("Nothing named yet.", 12.5, False, MUTED)])

        _rect(slide, Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.0), WHITE, NAVY)
        _text(slide, Inches(1.0), Inches(5.08), Inches(11.3), Inches(0.8),
              (f"{campus}: {cover.get('filled', 0)} students rated the week "
               f"{_fmt(head.get('vibe'), 1)}/10, belonging {_fmt(head.get('belonging'), 1)}/10 "
               f"and their chance of succeeding here {_fmt(head.get('success'), 1)}/10. "
               + (f"What stressed them most was "
                  f"{clean(_top(highlights.get('stressors')).get('label', ''), 40)}."
                  if highlights.get("stressors") else "")),
              size=12.5, colour=INK, spacing=1.3)

        # ── Department scoreboard ──────────────────────────────────────────────
        #
        # Every department that has a student in scope, over as many slides as
        # that takes. It used to be one slide sliced to the top twelve by vibe,
        # which meant the departments cut were by construction the lowest-rated
        # ones — the reader saw a leaderboard and was told it was a roster.
        if departments:
            ranked = sorted(
                departments,
                key=lambda r: (r.get("vibe") is None, -(r.get("vibe") or 0),
                               -r.get("filled", 0), r["dept"].lower()))
            small = sum(1 for r in ranked if not r.get("reportable", True))
            pages = [ranked[i:i + PER_PAGE] for i in range(0, len(ranked), PER_PAGE)]

            for page, rows in enumerate(pages, 1):
                slide = _blank(prs)
                _stage(slide)
                _title(slide,
                       "Department Scoreboard" +
                       (f" ({page} of {len(pages)})" if len(pages) > 1 else ""),
                       top=Inches(1.6), size=28)
                _table(slide, Inches(0.75), Inches(2.62), Inches(11.8),
                       ["Department", "Answered", "Rate", "Vibe /10",
                        "Belonging /10", "Bridge /5", "NPS"],
                       [[clean(r["dept"], 42),
                         f"{r.get('filled', 0)} / {r.get('eligible', 0)}",
                         _fmt(r.get("pct"), 0, "%"),
                         _fmt(r.get("vibe"), 1),
                         _fmt(r.get("belonging"), 1),
                         _fmt(r.get("bridge"), 1),
                         _fmt(r.get("nps"), 0)] for r in rows],
                       widths=[0.34, 0.12, 0.09, 0.11, 0.14, 0.11, 0.09])

                footnote = (
                    f"All {len(ranked)} departments with a student in scope, "
                    f"best-rated first; {cover.get('filled', 0):,} responses in total.")
                if small:
                    footnote += (
                        f"  A dash means fewer than {MIN_REPORTABLE} students in that "
                        "department answered: the count is shown, the scores are "
                        "withheld because an average of one or two replies identifies "
                        f"the students who gave them ({small} of {len(ranked)} here).")
                _text(slide, Inches(0.75), H - Inches(1.28), Inches(11.8), Inches(0.62),
                      footnote, size=10.5, colour=MUTED, spacing=1.2)

        # ── Slide 14 · Closing ─────────────────────────────────────────────────
        slide = _blank(prs)
        _stage(slide)
        _title(slide, "In one line", top=Inches(1.9), size=26)
        _text(slide, Inches(1.2), Inches(3.1), Inches(10.9), Inches(2.0),
              (f"{total} students at {campus} rated Deeksharambh "
               f"{_fmt(head.get('vibe'), 1)} out of 10 — {mood_word.lower()} — "
               f"with a Net Promoter Score of {_fmt(head.get('nps'), 0)} and a "
               f"belonging score of {_fmt(head.get('belonging'), 1)} out of 10."),
              size=22, bold=True, colour=NAVY, align=PP_ALIGN.CENTER, spacing=1.3)
        _text(slide, Inches(1.2), Inches(5.4), Inches(10.9), Inches(0.4),
              "Office of Academics · JAIN (Deemed-to-be University)",
              size=12.5, colour=MUTED, align=PP_ALIGN.CENTER)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
