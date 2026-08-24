"""
deeksharambh_brief_export.py — the "Student Experience & Orientation Impact
Analysis" brief, as a deck, a Word document, and a workbook.

Built to the shape of the report the Office of Academics already circulates
by hand: a cover, the three-step journey (baseline, Deeksharambh, post
survey), every department's registered-vs-submitted count, the programme's
theme and its nine sections, what students said about the week, what they
are asking for next, and a closing summary. Every number on it comes from
`deeksharambh_brief()` in orientation_data.py — nothing here is composed or
estimated, only laid out.

Its own palette, matching the printed brief: a warm cream ground, deep navy
ink, and three accents — teal, gold, coral — cycled one per journey step.
Georgia for headlines, Arial for everything else.
"""
from __future__ import annotations

import io

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.orientation_charts import clean

# python-docx backs the Word half of this module only. Importing it at module
# load time meant a server without it installed couldn't even build the
# slide deck or the workbook — one missing optional dependency taking down
# two unrelated downloads. It loads lazily, the first time a docx is built.
try:
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches as DocxInches
    from docx.shared import Pt as DocxPt
    from docx.shared import RGBColor as DocxRGBColor
except ImportError:
    Document = WD_TABLE_ALIGNMENT = WD_ALIGN_PARAGRAPH = None
    OxmlElement = qn = DocxInches = DocxPt = None

    def DocxRGBColor(*_a, **_kw):  # placeholder so the module-level palette below still loads
        return None

# ── Palette ───────────────────────────────────────────────────────────────
CREAM = RGBColor(0xF7, 0xF4, 0xEF)
NAVY = RGBColor(0x0B, 0x1D, 0x36)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
CORAL = RGBColor(0xE0, 0x7A, 0x5F)
SLATE = RGBColor(0x64, 0x74, 0x8B)
INK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_TINT = RGBColor(0xF0, 0xED, 0xE6)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

STEP_TONES = [NAVY, TEAL, GOLD, CORAL]

DOCX_NAVY = DocxRGBColor(0x0B, 0x1D, 0x36)
DOCX_TEAL = DocxRGBColor(0x0F, 0x76, 0x6E)
DOCX_GOLD = DocxRGBColor(0xD4, 0xA0, 0x17)
DOCX_CORAL = DocxRGBColor(0xE0, 0x7A, 0x5F)
DOCX_SLATE = DocxRGBColor(0x64, 0x74, 0x8B)
DOCX_INK = DocxRGBColor(0x1E, 0x29, 0x3B)
DOCX_WHITE = DocxRGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"
SANS = "Arial"

# The printed brief's own page size — 10 x 5.625in, not the 13.333in the
# other decks use.
W = Inches(10)
H = Inches(5.625)
RECT = 1  # MSO_SHAPE.RECTANGLE

# The programme's own stated theme — printed on every Deeksharambh 2026
# cover, not something this brief derives from a count.
THEME_TITLE = "Human + AI: Building Your Future"
THEME_LINE = (
    "Smooth transition into campus life  ·  Bridge Course for readiness  ·  "
    "Immersive orientation that builds belonging & confidence"
)

SECTION_TAGLINES = {
    "The Vibe Check": "Overall feel of the first week",
    "Settling In": "Ease of transition to campus life",
    "Footsteps (pre-arrival)": "Pre-arrival guidance module",
    "Orientation Experience": "Core sessions & campus immersion",
    "Bridge Course": "Academic foundation & readiness",
    "NEP 2020 & Digital Readiness": "Policy alignment & digital skills",
    "The Gen Z Lens": "Student voice & generational view",
    "Belonging & Expectations": "Sense of community & future goals",
    "Score & Mic Drop": "Key metrics & final takeaways",
}


def _sections() -> list[tuple[str, str]]:
    """The report's own nine sections, titled and tagged — not invented here."""
    from app.orientation_analysis import SECTIONS

    return [(clean(title), SECTION_TAGLINES.get(clean(title), "")) for title, _ in SECTIONS]


# ── PPTX ──────────────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    return slide


def _text(slide, left, top, width, height, text, *, size=11, bold=False,
          colour=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          font=SANS, spacing=0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        for f in [para.font] + [r.font for r in para.runs]:
            f.size = Pt(size)
            f.bold = bold
            f.color.rgb = colour
            f.name = font
        if spacing:
            for run in para.runs or [para.add_run()]:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(spacing * 100)))
    return box


def _rect(slide, left, top, width, height, *, fill=WHITE, line=None):
    box = slide.shapes.add_shape(RECT, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.5)
    box.shadow.inherit = False
    return box


def _kicker(slide, left, top, text, *, colour=TEAL, width=Inches(9)):
    _text(slide, left, top, width, Inches(0.25), text.upper(),
          size=11, bold=True, colour=colour, spacing=1.6)


def _footer(slide, generated_at: str, campus: str):
    _text(slide, Inches(0.45), H - Inches(0.35), Inches(9), Inches(0.25),
          f"Deeksharambh 2026 · {campus} · generated {generated_at}",
          size=7.5, colour=SLATE)


def _journey_card(slide, left, top, width, height, step: int, tone,
                   count: str, label: str, sub: str):
    _rect(slide, left, top, width, height, fill=WHITE)
    _rect(slide, left, top, width, Inches(0.06), fill=tone)
    _text(slide, left + Inches(0.12), top + Inches(0.14), width - Inches(0.24), Inches(0.2),
          f"0{step}", size=10, colour=SLATE)
    _text(slide, left + Inches(0.12), top + Inches(0.32), width - Inches(0.24), Inches(0.42),
          count, size=26, bold=True, colour=NAVY)
    _text(slide, left + Inches(0.12), top + height - Inches(0.62), width - Inches(0.24), Inches(0.22),
          label, size=12.5, bold=True, colour=INK)
    _text(slide, left + Inches(0.12), top + height - Inches(0.38), width - Inches(0.24), Inches(0.3),
          sub, size=10, colour=SLATE)


def _slide_cover(prs, data, generated_at):
    slide = _blank(prs)
    _kicker(slide, Inches(0.6), Inches(1.35), "DEEKSHARAMBH 2026")
    _text(slide, Inches(0.55), Inches(1.7), Inches(8.9), Inches(1.3),
          "Student Experience &\nOrientation Impact Analysis",
          size=30, bold=True, colour=NAVY, font=SERIF)
    _text(slide, Inches(0.6), Inches(3.05), Inches(8.8), Inches(0.3),
          f"{data['campus']}  ·  All Departments  ·  UG & PG", size=11, colour=SLATE)
    _text(slide, Inches(0.6), Inches(3.35), Inches(8.8), Inches(0.25),
          "Source: student-survey.juooa.cloud/shared/orientation", size=9, colour=SLATE)
    _text(slide, Inches(0.6), H - Inches(0.55), Inches(8.8), Inches(0.3),
          f"Office of Academics  ·  JAIN (Deemed-to-be University)  ·  {generated_at}",
          size=9, colour=SLATE)


def _slide_journey(prs, data, generated_at):
    slide = _blank(prs)
    _kicker(slide, Inches(0.35), Inches(0.18), "01  ·  STUDENT JOURNEY")
    _text(slide, Inches(0.35), Inches(0.42), Inches(9), Inches(0.35),
          "Baseline · Deeksharambh · Post survey", size=18, bold=True, colour=NAVY, font=SERIF)

    stages = data["journey"]["stages"]
    labels = [("Registered", "Total cohort"), ("Baseline", "Completed"),
              ("Deeksharambh", "Submitted"), ("Post survey", "Completed")]
    card_w = Inches(2.2)
    gap = Inches(0.16)
    for i, stage in enumerate(stages):
        left = Inches(0.35) + i * (card_w + gap)
        label, sub = labels[i]
        _journey_card(slide, left, Inches(0.95), card_w, Inches(1.35),
                      i + 1, STEP_TONES[i], f"{stage['count']:,}", label, sub)

    _rect(slide, Inches(0.35), Inches(2.55), Inches(9.3), Inches(1.0), fill=NAVY)
    _text(slide, Inches(0.55), Inches(2.68), Inches(8.9), Inches(0.22),
          "FULL JOURNEY COMPLETED", size=10.5, bold=True, colour=GOLD, spacing=1)
    reg = data["journey"]["registered"]
    _text(slide, Inches(0.55), Inches(2.92), Inches(8.9), Inches(0.28),
          f"{data['full_journey']:,} students ({data['full_journey_pct']:.0f}%) completed all three "
          f"steps — Baseline → Deeksharambh → Post survey", size=12, colour=WHITE)
    _text(slide, Inches(0.55), Inches(3.2), Inches(8.9), Inches(0.28),
          f"{data['pre_not_deeksharambh']:,} students finished the baseline but have not yet "
          f"submitted the Deeksharambh form", size=10.5, colour=RGBColor(0xE2, 0xE8, 0xF0))

    _footer(slide, generated_at, data["campus"])


def _slide_departments(prs, rows, generated_at, campus, part, of):
    slide = _blank(prs)
    suffix = f" — part {part} of {of}" if of > 1 else ""
    _kicker(slide, Inches(0.35), Inches(0.15), f"02  ·  ALL DEPARTMENTS{suffix}")
    _text(slide, Inches(0.35), Inches(0.37), Inches(9), Inches(0.32),
          "Registered  ·  Submitted Deeksharambh Survey", size=15, bold=True, colour=NAVY, font=SERIF)

    _rect(slide, Inches(0.3), Inches(0.8), Inches(9.4), Inches(0.38), fill=NAVY)
    _text(slide, Inches(0.45), Inches(0.87), Inches(5.5), Inches(0.25), "Department",
          size=11, bold=True, colour=WHITE)
    _text(slide, Inches(6.1), Inches(0.87), Inches(1.6), Inches(0.25), "Registered",
          size=11, bold=True, colour=WHITE, align=PP_ALIGN.RIGHT)
    _text(slide, Inches(7.85), Inches(0.87), Inches(1.6), Inches(0.25), "Submitted",
          size=11, bold=True, colour=WHITE, align=PP_ALIGN.RIGHT)

    row_h = Inches(0.135)
    top = Inches(1.22)
    for i, row in enumerate(rows):
        if i % 2 == 0:
            _rect(slide, Inches(0.3), top, Inches(9.4), row_h, fill=ROW_TINT)
        _text(slide, Inches(0.45), top, Inches(5.5), row_h, clean(row["dept"], 60),
              size=9.5, colour=INK)
        _text(slide, Inches(6.1), top, Inches(1.6), row_h, str(row["eligible"]),
              size=9.5, colour=INK, align=PP_ALIGN.RIGHT)
        _text(slide, Inches(7.85), top, Inches(1.6), row_h, str(row["filled"]),
              size=9.5, colour=INK, align=PP_ALIGN.RIGHT)
        top += row_h

    _footer(slide, generated_at, campus)


def _slide_theme(prs, generated_at, campus):
    slide = _blank(prs)
    _kicker(slide, Inches(0.45), Inches(0.18), "03  ·  THEME & SECTIONS")
    _text(slide, Inches(0.45), Inches(0.4), Inches(9), Inches(0.4),
          THEME_TITLE, size=19, bold=True, colour=NAVY, font=SERIF)
    _rect(slide, Inches(0.4), Inches(0.9), Inches(9.2), Inches(0.55), fill=NAVY)
    _text(slide, Inches(0.55), Inches(1.0), Inches(8.9), Inches(0.35),
          f"Theme: {THEME_LINE}", size=10, colour=WHITE)

    sections = _sections()
    cols, col_w, row_h = 3, Inches(3.0), Inches(1.05)
    gap_x, gap_y = Inches(0.15), Inches(0.15)
    for i, (title, tagline) in enumerate(sections):
        r, c = divmod(i, cols)
        left = Inches(0.4) + c * (col_w + gap_x)
        top = Inches(1.7) + r * (row_h + gap_y)
        _rect(slide, left, top, col_w, row_h, fill=WHITE)
        _rect(slide, left + Inches(0.15), top + Inches(0.2), Inches(0.4), Inches(0.3), fill=TEAL)
        _text(slide, left + Inches(0.15), top + Inches(0.22), Inches(0.4), Inches(0.28),
              f"{i + 1:02d}", size=11, colour=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, left + Inches(0.65), top + Inches(0.18), col_w - Inches(0.8), Inches(0.3),
              title, size=12, bold=True, colour=NAVY)
        _text(slide, left + Inches(0.65), top + Inches(0.52), col_w - Inches(0.8), Inches(0.35),
              tagline, size=9.5, colour=SLATE)

    _footer(slide, generated_at, campus)


def _stat_panel(slide, left, top, width, height, value: str, label: str):
    _rect(slide, left, top, width, height, fill=WHITE)
    _text(slide, left, top + Inches(0.2), width, Inches(0.6), value,
          size=34, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)
    _text(slide, left + Inches(0.2), top + Inches(0.9), width - Inches(0.4), Inches(0.6),
          label, size=11, colour=SLATE, align=PP_ALIGN.CENTER)


def _slide_overall(prs, data, generated_at):
    slide = _blank(prs)
    _kicker(slide, Inches(0.45), Inches(0.2), "04  ·  OVERALL DEEKSHARAMBH")
    _text(slide, Inches(0.45), Inches(0.45), Inches(9), Inches(0.35),
          "Who Took It & What They Said", size=18, bold=True, colour=NAVY, font=SERIF)

    _stat_panel(slide, Inches(0.4), Inches(1.0), Inches(4.5), Inches(1.8),
               f"{data['deeksharambh_count']:,}",
               "students submitted the Deeksharambh survey")
    _stat_panel(slide, Inches(5.15), Inches(1.0), Inches(4.45), Inches(1.8),
               f"{data['deeksharambh_pct_of_pre']:.0f}%",
               f"of the {data['pre_completed']:,} who completed the baseline took Deeksharambh")

    _rect(slide, Inches(0.4), Inches(3.05), Inches(9.2), Inches(2.15), fill=WHITE)
    _text(slide, Inches(0.65), Inches(3.2), Inches(8.7), Inches(0.3),
          "What students said about the week", size=13, bold=True, colour=NAVY)

    rows = []
    if data.get("welcomed"):
        rows.append((data["welcomed"]["pct"], f"Felt welcomed — {clean(data['welcomed']['label'])}"))
    if data.get("transition"):
        rows.append((data["transition"]["pct"], f"Transition into university life — {clean(data['transition']['label'])}"))
    if data.get("help_contact"):
        rows.append((data["help_contact"]["pct"], f"Knows whom to ask for help — {clean(data['help_contact']['label'])}"))
    rows.append((data["high_vibe_pct"], "Rated the week 8 or higher out of 10"))

    positions = [(0.7, 3.65), (5.2, 3.65), (0.7, 4.3), (5.2, 4.3)]
    for (pct, label), (x, y) in zip(rows, positions):
        _text(slide, Inches(x), Inches(y), Inches(0.8), Inches(0.35), f"{pct:.0f}%",
              size=17, bold=True, colour=TEAL)
        _text(slide, Inches(x + 0.9), Inches(y), Inches(3.4), Inches(0.35), label,
              size=10.5, colour=INK)

    _footer(slide, generated_at, data["campus"])


def _slide_expectations(prs, data, generated_at):
    slide = _blank(prs)
    _kicker(slide, Inches(0.45), Inches(0.2), "05  ·  STUDENT EXPECTATIONS")
    _text(slide, Inches(0.45), Inches(0.45), Inches(9), Inches(0.35),
          "What Students Are Asking For Next", size=18, bold=True, colour=NAVY, font=SERIF)

    items = data.get("expectations") or []
    card_w, card_h = Inches(4.55), Inches(0.9)
    gap_x, gap_y = Inches(0.2), Inches(0.15)
    for i, opt in enumerate(items[:8]):
        r, c = divmod(i, 2)
        left = Inches(0.4) + c * (card_w + gap_x)
        top = Inches(1.0) + r * (card_h + gap_y)
        _rect(slide, left, top, card_w, card_h, fill=WHITE)
        _text(slide, left + Inches(0.2), top + Inches(0.13), card_w - Inches(0.4), Inches(0.3),
              clean(opt["label"], 46), size=12, bold=True, colour=NAVY)
        _text(slide, left + Inches(0.2), top + Inches(0.46), card_w - Inches(0.4), Inches(0.35),
              f"{opt['count']} students · {opt['pct']:.0f}% of those who answered",
              size=9.5, colour=SLATE)

    _footer(slide, generated_at, data["campus"])


def _slide_conclusion(prs, data, generated_at):
    slide = _blank(prs)
    _rect(slide, Emu(0), Emu(0), Inches(0.18), H, fill=GOLD)
    _text(slide, Inches(0.7), Inches(0.55), Inches(8.5), Inches(0.35),
          "CONCLUSION & SUMMARY", size=13, bold=True, colour=TEAL, spacing=1)

    welcomed_bit = clean(data["welcomed"]["label"]).lower() if data.get("welcomed") else "welcomed"
    transition_bit = clean(data["transition"]["label"]).lower() if data.get("transition") else "smooth"
    para1 = (
        f"{data['deeksharambh_count']:,} students submitted the Deeksharambh survey. "
        f"Of those who rated the week, most said they felt {welcomed_bit}, found the transition "
        f"{transition_bit}, and {data['high_vibe_pct']:.0f}% rated the week 8 or higher out of 10."
    )
    _text(slide, Inches(0.7), Inches(1.0), Inches(8.5), Inches(0.85), para1, size=13, colour=INK)

    landed = [clean(s["label"]) for s in data.get("sessions_landed") or []]
    weak = [clean(s["label"]) for s in data.get("needs_work") or []]
    bits = []
    if landed:
        bits.append(f"{', '.join(landed)} had the strongest impact")
    if weak:
        bits.append(f"{', '.join(weak)} remain the main areas to strengthen")
    para2 = ". ".join(bits) + "." if bits else ""
    if para2:
        _text(slide, Inches(0.7), Inches(1.9), Inches(8.5), Inches(0.85), para2, size=13, colour=INK)

    _rect(slide, Inches(0.7), Inches(2.85), Inches(8.5), Inches(1.05), fill=NAVY)
    _text(slide, Inches(0.7), Inches(3.03), Inches(8.5), Inches(0.35), "THANK YOU",
          size=20, bold=True, colour=WHITE, align=PP_ALIGN.CENTER, font=SERIF)
    _text(slide, Inches(0.7), Inches(3.5), Inches(8.5), Inches(0.25),
          "Office of Academics  ·  JAIN (Deemed-to-be University)",
          size=10, colour=RGBColor(0xE2, 0xE8, 0xF0), align=PP_ALIGN.CENTER)
    _footer(slide, generated_at, data["campus"])


def build_deeksharambh_brief_pptx(data: dict, *, generated_at: str) -> bytes:
    """The seven-slide brief: cover, journey, departments, theme, overall,
    expectations, conclusion — the shape of the report already circulated."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    _slide_cover(prs, data, generated_at)
    _slide_journey(prs, data, generated_at)

    PAGE_SIZE = 26
    depts = data["departments"]
    pages = [depts[i:i + PAGE_SIZE] for i in range(0, len(depts), PAGE_SIZE)] or [[]]
    for part, rows in enumerate(pages, start=1):
        _slide_departments(prs, rows, generated_at, data["campus"], part, len(pages))

    _slide_theme(prs, generated_at, data["campus"])
    _slide_overall(prs, data, generated_at)
    _slide_expectations(prs, data, generated_at)
    _slide_conclusion(prs, data, generated_at)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ── DOCX ──────────────────────────────────────────────────────────────────

def _shade(cell, hex_colour: str) -> None:
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_colour)
    cell._tc.get_or_add_tcPr().append(shd)


def _no_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "nil")
        borders.append(el)
    tbl_pr.append(borders)


def _run(paragraph, text, *, size=11, bold=False, colour=DOCX_INK, align=None):
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run(text)
    run.font.size = DocxPt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return run


def _cell_text(cell, text, **kw):
    cell.text = ""
    return _run(cell.paragraphs[0], text, **kw)


def _docx_heading(doc: Document, kicker: str, title: str) -> None:
    p = doc.add_paragraph()
    _run(p, kicker.upper(), size=10, bold=True, colour=DOCX_TEAL)
    h = doc.add_heading(title, level=1)
    for run in h.runs:
        run.font.color.rgb = DOCX_NAVY
        run.font.name = SERIF


def _docx_cover(doc: Document, data: dict, generated_at: str) -> None:
    p = doc.add_paragraph()
    _run(p, "DEEKSHARAMBH 2026", size=12, bold=True, colour=DOCX_TEAL)
    h = doc.add_heading("Student Experience & Orientation Impact Analysis", level=0)
    for run in h.runs:
        run.font.color.rgb = DOCX_NAVY
        run.font.name = SERIF
    p = doc.add_paragraph()
    _run(p, f"{data['campus']}  ·  All Departments  ·  UG & PG", size=11, colour=DOCX_SLATE)
    p = doc.add_paragraph()
    _run(p, "Source: student-survey.juooa.cloud/shared/orientation", size=9, colour=DOCX_SLATE)
    p = doc.add_paragraph()
    _run(p, f"Office of Academics · JAIN (Deemed-to-be University) · {generated_at}",
         size=9, colour=DOCX_SLATE)
    doc.add_page_break()


def _docx_journey(doc: Document, data: dict) -> None:
    _docx_heading(doc, "01 · Student Journey", "Baseline · Deeksharambh · Post survey")
    stages = data["journey"]["stages"]
    labels = ["Registered", "Baseline", "Deeksharambh", "Post survey"]
    table = doc.add_table(rows=1, cols=4)
    table.style = "Light Grid Accent 1"
    for c, (label, stage) in enumerate(zip(labels, stages)):
        _cell_text(table.rows[0].cells[c], f"{label}\n{stage['count']:,}", bold=True, size=12)
    doc.add_paragraph()
    p = doc.add_paragraph()
    _run(p, "Full journey completed: ", bold=True)
    _run(p, f"{data['full_journey']:,} students ({data['full_journey_pct']:.0f}%) completed all "
            f"three steps — Baseline → Deeksharambh → Post survey.")
    p = doc.add_paragraph()
    _run(p, f"{data['pre_not_deeksharambh']:,} students finished the baseline but have not yet "
            f"submitted the Deeksharambh form.", colour=DOCX_SLATE)
    doc.add_paragraph()


def _docx_departments(doc: Document, data: dict) -> None:
    _docx_heading(doc, "02 · All Departments", "Registered · Submitted Deeksharambh Survey")
    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    widths = (DocxInches(4.2), DocxInches(1.3), DocxInches(1.3))
    for c, (label, width) in enumerate(zip(("Department", "Registered", "Submitted"), widths)):
        cell = table.rows[0].cells[c]
        cell.width = width
        _shade(cell, "0B1D36")
        align = WD_ALIGN_PARAGRAPH.RIGHT if c else None
        _cell_text(cell, label, bold=True, colour=DOCX_WHITE, align=align)
    for i, row in enumerate(data["departments"]):
        cells = table.add_row().cells
        for c, width in enumerate(widths):
            cells[c].width = width
        _cell_text(cells[0], clean(row["dept"], 60), size=10)
        _cell_text(cells[1], str(row["eligible"]), size=10, align=WD_ALIGN_PARAGRAPH.RIGHT)
        _cell_text(cells[2], str(row["filled"]), size=10, align=WD_ALIGN_PARAGRAPH.RIGHT)
        if i % 2 == 0:
            for cell in cells:
                _shade(cell, "F0EDE6")
    doc.add_paragraph()


def _docx_theme(doc: Document) -> None:
    _docx_heading(doc, "03 · Theme & Sections", THEME_TITLE)
    p = doc.add_paragraph()
    _run(p, f"Theme: {THEME_LINE}", size=10.5, colour=DOCX_SLATE)
    doc.add_paragraph()
    for i, (title, tagline) in enumerate(_sections(), start=1):
        p = doc.add_paragraph()
        _run(p, f"{i:02d}  ", bold=True, colour=DOCX_TEAL)
        _run(p, title, bold=True, colour=DOCX_NAVY)
        _run(p, f"  —  {tagline}", colour=DOCX_SLATE, size=10)
    doc.add_paragraph()


def _docx_overall(doc: Document, data: dict) -> None:
    _docx_heading(doc, "04 · Overall Deeksharambh", "Who Took It & What They Said")
    table = doc.add_table(rows=1, cols=2)
    table.style = "Light Grid Accent 1"
    _cell_text(table.rows[0].cells[0],
               f"{data['deeksharambh_count']:,}\nstudents submitted the Deeksharambh survey",
               bold=True, size=12)
    _cell_text(table.rows[0].cells[1],
               f"{data['deeksharambh_pct_of_pre']:.0f}%\nof the {data['pre_completed']:,} who "
               f"completed the baseline took Deeksharambh", bold=True, size=12)
    doc.add_paragraph()
    p = doc.add_paragraph()
    _run(p, "What students said about the week", bold=True, colour=DOCX_NAVY)

    rows = []
    if data.get("welcomed"):
        rows.append((data["welcomed"]["pct"], f"Felt welcomed — {clean(data['welcomed']['label'])}"))
    if data.get("transition"):
        rows.append((data["transition"]["pct"], f"Transition into university life — {clean(data['transition']['label'])}"))
    if data.get("help_contact"):
        rows.append((data["help_contact"]["pct"], f"Knows whom to ask for help — {clean(data['help_contact']['label'])}"))
    rows.append((data["high_vibe_pct"], "Rated the week 8 or higher out of 10"))
    for pct, label in rows:
        p = doc.add_paragraph(style="List Bullet")
        _run(p, f"{pct:.0f}%  ", bold=True, colour=DOCX_TEAL)
        _run(p, label)
    doc.add_paragraph()


def _docx_expectations(doc: Document, data: dict) -> None:
    _docx_heading(doc, "05 · Student Expectations", "What Students Are Asking For Next")
    for opt in (data.get("expectations") or [])[:8]:
        p = doc.add_paragraph(style="List Bullet")
        _run(p, clean(opt["label"], 60), bold=True, colour=DOCX_NAVY)
        _run(p, f"  —  {opt['count']} students · {opt['pct']:.0f}% of those who answered",
             colour=DOCX_SLATE, size=10)
    doc.add_paragraph()


def _docx_conclusion(doc: Document, data: dict) -> None:
    _docx_heading(doc, "Conclusion & Summary", "In one page")
    welcomed_bit = clean(data["welcomed"]["label"]).lower() if data.get("welcomed") else "welcomed"
    transition_bit = clean(data["transition"]["label"]).lower() if data.get("transition") else "smooth"
    p = doc.add_paragraph()
    _run(p, f"{data['deeksharambh_count']:,} students submitted the Deeksharambh survey. "
            f"Of those who rated the week, most said they felt {welcomed_bit}, found the "
            f"transition {transition_bit}, and {data['high_vibe_pct']:.0f}% rated the week 8 "
            f"or higher out of 10.")
    landed = [clean(s["label"]) for s in data.get("sessions_landed") or []]
    weak = [clean(s["label"]) for s in data.get("needs_work") or []]
    bits = []
    if landed:
        bits.append(f"{', '.join(landed)} had the strongest impact")
    if weak:
        bits.append(f"{', '.join(weak)} remain the main areas to strengthen")
    if bits:
        p = doc.add_paragraph()
        _run(p, ". ".join(bits) + ".")


def build_deeksharambh_brief_docx(data: dict, *, generated_at: str) -> bytes:
    """The same brief, laid out as a Word document."""
    if Document is None:
        raise RuntimeError(
            "python-docx is not installed — run `pip install -r requirements.txt`")
    doc = Document()
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    for margin in ("left_margin", "right_margin", "top_margin", "bottom_margin"):
        setattr(section, margin, DocxInches(0.9))

    style = doc.styles["Normal"]
    style.font.name = SANS
    style.font.size = DocxPt(11)
    style.font.color.rgb = DOCX_INK

    _docx_cover(doc, data, generated_at)
    _docx_journey(doc, data)
    _docx_departments(doc, data)
    doc.add_page_break()
    _docx_theme(doc)
    _docx_overall(doc, data)
    _docx_expectations(doc, data)
    doc.add_page_break()
    _docx_conclusion(doc, data)

    footer = doc.sections[0].footer
    footer.paragraphs[0].text = f"Deeksharambh 2026 · {data['campus']} · generated {generated_at}"
    footer.paragraphs[0].runs[0].font.size = DocxPt(8.5)
    footer.paragraphs[0].runs[0].font.color.rgb = DOCX_SLATE

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


# ── XLSX ──────────────────────────────────────────────────────────────────

def build_deeksharambh_brief_xlsx(data: dict, *, generated_at: str) -> bytes:
    """The same brief's numbers, as a workbook — one sheet per section, so
    the department table can be filtered and re-sorted like any roster."""
    wb = Workbook()

    header_fill = PatternFill("solid", fgColor="0B1D36")
    tint_fill = PatternFill("solid", fgColor="F0EDE6")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    title_font = Font(bold=True, color="0B1D36", size=14)
    muted_font = Font(color="64748B", size=10)
    thin = Side(style="thin", color="E2E8F0")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    centre = Alignment(horizontal="center", vertical="center")

    def sized(ws, widths):
        for i, w in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w

    # ── Overview ─────────────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Overview"
    ws["A1"] = "Deeksharambh 2026 · Student Experience & Orientation Impact Analysis"
    ws["A1"].font = title_font
    ws["A2"] = f"{data['campus']} · generated {generated_at}"
    ws["A2"].font = muted_font
    sized(ws, [34, 16])

    rows = [
        ("Registered", data["journey"]["registered"]),
        ("Completed the baseline", data["journey"]["stages"][1]["count"]),
        ("Submitted Deeksharambh", data["deeksharambh_count"]),
        ("Completed the post survey", data["journey"]["stages"][3]["count"]),
        ("Completed all three steps", data["full_journey"]),
        ("Completed all three steps (%)", data["full_journey_pct"]),
        ("Finished baseline, not yet Deeksharambh", data["pre_not_deeksharambh"]),
        ("Deeksharambh as % of baseline completers", data["deeksharambh_pct_of_pre"]),
        ("Rated the week 8+ / 10 (%)", data["high_vibe_pct"]),
    ]
    r = 4
    for label, value in rows:
        ws.cell(r, 1, label).font = Font(bold=True)
        ws.cell(r, 2, value)
        r += 1

    # ── Departments ──────────────────────────────────────────────────────
    ws2 = wb.create_sheet("Departments")
    ws2.append(["Department", "Registered", "Submitted", "Response rate %"])
    for cell in ws2[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = centre
        cell.border = border
    for i, row in enumerate(data["departments"], start=2):
        ws2.cell(i, 1, clean(row["dept"], 80))
        ws2.cell(i, 2, row["eligible"])
        ws2.cell(i, 3, row["filled"])
        ws2.cell(i, 4, row["pct"])
        if i % 2 == 0:
            for c in range(1, 5):
                ws2.cell(i, c).fill = tint_fill
        for c in range(1, 5):
            ws2.cell(i, c).border = border
    sized(ws2, [42, 13, 13, 15])
    ws2.freeze_panes = "A2"

    # ── What students said ───────────────────────────────────────────────
    ws3 = wb.create_sheet("What students said")
    ws3.append(["Question", "Top answer", "%"])
    for cell in ws3[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = centre
        cell.border = border
    said = []
    if data.get("welcomed"):
        said.append(("Felt welcomed", clean(data["welcomed"]["label"]), data["welcomed"]["pct"]))
    if data.get("transition"):
        said.append(("Transition into university life", clean(data["transition"]["label"]),
                     data["transition"]["pct"]))
    if data.get("help_contact"):
        said.append(("Knows whom to ask for help", clean(data["help_contact"]["label"]),
                     data["help_contact"]["pct"]))
    said.append(("Rated the week 8 or higher / 10", "", data["high_vibe_pct"]))
    for i, (q, a, pct) in enumerate(said, start=2):
        ws3.cell(i, 1, q)
        ws3.cell(i, 2, a)
        ws3.cell(i, 3, pct)
        for c in range(1, 4):
            ws3.cell(i, c).border = border
    sized(ws3, [34, 34, 10])

    # ── Expectations ─────────────────────────────────────────────────────
    ws4 = wb.create_sheet("Expectations")
    ws4.append(["What students are asking for", "Students", "% of those who answered"])
    for cell in ws4[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = centre
        cell.border = border
    for i, opt in enumerate(data.get("expectations") or [], start=2):
        ws4.cell(i, 1, clean(opt["label"]))
        ws4.cell(i, 2, opt["count"])
        ws4.cell(i, 3, opt["pct"])
        for c in range(1, 4):
            ws4.cell(i, c).border = border
    sized(ws4, [42, 12, 20])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()
