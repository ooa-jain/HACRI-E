"""
PowerPoint export helper using python-pptx.
Generates slides with department stats and visual charts.
"""
from __future__ import annotations
import io
import os
import tempfile
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.charts import plot_cohort_png, plot_histograms_png, plot_h1_histogram_custom

def generate_dept_ppt(
    dept_name: str,
    users_list: list[dict],
    matched_data: dict[str, dict]
) -> bytes:
    # 1. Initialize presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme colors
    c_navy = RGBColor(27, 42, 74)       # #1B2A4A
    c_gold = RGBColor(201, 168, 76)     # #C9A84C
    c_teal = RGBColor(13, 148, 136)     # #0D9488
    c_dark = RGBColor(30, 41, 59)       # #1E293B
    c_light = RGBColor(245, 246, 250)   # #F5F6FA

    # Count stats
    total = len(users_list)
    pre_done = sum(1 for u in users_list if u.get("status") in ("pre_done", "post_done"))
    post_done = sum(1 for u in users_list if u.get("status") == "post_done")
    pending = pre_done - post_done

    # Slide 1: Title Slide (Dark Background)
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background fill (Navy)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = c_navy

    # Title Text
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HACRI-E2 Department Analysis"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = c_gold
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = f"Department: {dept_name or 'All Departments'}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = f"Generated on: {datetime.now().strftime('%d %B %Y')}"
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(180, 180, 180)
    p3.alignment = PP_ALIGN.LEFT

    # Slide 2: Statistics Summary
    slide2 = prs.slides.add_slide(slide_layout)
    
    # Title
    t_box = slide2.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    t_tf = t_box.text_frame
    t_p = t_tf.paragraphs[0]
    t_p.text = "Cohort Enrollment & Participation"
    t_p.font.size = Pt(28)
    t_p.font.bold = True
    t_p.font.color.rgb = c_navy
    
    # Draw stats cards
    stats = [
        ("Registered Students", str(total), c_navy),
        ("Pre-Survey Done", str(pre_done), c_teal),
        ("Post-Survey Done", str(post_done), c_gold),
        ("Pending Post-Survey", str(pending), RGBColor(220, 38, 38)) # Red
    ]
    
    card_width = Inches(2.6)
    card_height = Inches(2.8)
    gap = Inches(0.4)
    start_left = Inches(0.75)
    start_top = Inches(2.2)
    
    for idx, (label, val, val_color) in enumerate(stats):
        left = start_left + idx * (card_width + gap)
        # Create shape (rectangle)
        shape = slide2.shapes.add_shape(
            1, # Rectangle
            left, start_top, card_width, card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_light
        shape.line.color.rgb = RGBColor(220, 220, 220)
        
        # Add text to card
        card_tf = shape.text_frame
        card_tf.word_wrap = True
        card_tf.margin_top = Inches(0.4)
        
        p_val = card_tf.paragraphs[0]
        p_val.text = val
        p_val.font.size = Pt(48)
        p_val.font.bold = True
        p_val.font.color.rgb = val_color
        p_val.alignment = PP_ALIGN.CENTER
        
        p_lbl = card_tf.add_paragraph()
        p_lbl.text = f"\n{label}"
        p_lbl.font.size = Pt(14)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = c_dark
        p_lbl.alignment = PP_ALIGN.CENTER

    # Generate charts for slides
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        cohort_img = tmp_path / "cohort.png"
        hist_img = tmp_path / "histograms.png"
        h1_img = tmp_path / "h1.png"
        
        # Plot Matplotlib charts
        plot_cohort_png(matched_data, out_path=cohort_img)
        plot_histograms_png(matched_data, out_path=hist_img)
        plot_h1_histogram_custom(matched_data, out_path=h1_img)

        # Slide 3: Quadrant Score Chart
        if cohort_img.exists():
            slide3 = prs.slides.add_slide(slide_layout)
            # Title
            t3 = slide3.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
            t3_tf = t3.text_frame
            t3_p = t3_tf.paragraphs[0]
            t3_p.text = "AI Literacy × AI Readiness Quadrant"
            t3_p.font.size = Pt(28)
            t3_p.font.bold = True
            t3_p.font.color.rgb = c_navy
            
            # Insert Chart Image
            slide3.shapes.add_picture(
                str(cohort_img),
                left=Inches(3.166), top=Inches(1.5),
                width=Inches(7.0), height=Inches(5.2)
            )

        # Slide 4: Histograms Score Distributions Chart
        if hist_img.exists():
            slide4 = prs.slides.add_slide(slide_layout)
            # Title
            t4 = slide4.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
            t4_tf = t4.text_frame
            t4_p = t4_tf.paragraphs[0]
            t4_p.text = "Pre & Post Score Distributions"
            t4_p.font.size = Pt(28)
            t4_p.font.bold = True
            t4_p.font.color.rgb = c_navy
            
            # Insert Chart Image
            slide4.shapes.add_picture(
                str(hist_img),
                left=Inches(1.666), top=Inches(1.6),
                width=Inches(10.0), height=Inches(4.8)
            )

        # Slide 5: Section H1 Chart
        if h1_img.exists():
            slide5 = prs.slides.add_slide(slide_layout)
            # Title
            t5 = slide5.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
            t5_tf = t5.text_frame
            t5_p = t5_tf.paragraphs[0]
            t5_p.text = "H1 · Post-Workshop Understanding Change"
            t5_p.font.size = Pt(28)
            t5_p.font.bold = True
            t5_p.font.color.rgb = c_navy
            
            # Insert Chart Image
            slide5.shapes.add_picture(
                str(h1_img),
                left=Inches(2.166), top=Inches(1.8),
                width=Inches(9.0), height=Inches(4.5)
            )

    # 4. Save presentation to buffer
    out_buf = io.BytesIO()
    prs.save(out_buf)
    return out_buf.getvalue()
