"""
cohort_ppt.py — the outcome and impact report as a slide deck.

Same numbers as the dashboard page, laid out for a committee room: the journey
first, then the baseline, then what changed, then the department detail.
"""
from __future__ import annotations

import io
import tempfile
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.cohort_charts import (
    clean, plot_campus_split, plot_departments, plot_distribution, plot_journey,
    plot_mix, plot_movement, plot_pre_post,
)

NAVY = RGBColor(0x0D, 0x21, 0x47)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x64, 0x74, 0x8B)
MIST = RGBColor(0xF1, 0xF5, 0xF9)
TEAL = RGBColor(0x0D, 0x94, 0x88)
ROSE = RGBColor(0xBE, 0x12, 0x3C)
DEEP = RGBColor(0x1E, 0x3A, 0x8A)

W = Inches(13.333)
H = Inches(7.5)
RECT = 1
ROUND = 5


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, left, top, width, height, text, *, size=18, bold=False,
          colour=INK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.text = text
    para.alignment = align
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.italic = italic
    para.font.color.rgb = colour
    return box


def _heading(slide, title: str, kicker: str = "") -> None:
    bar = slide.shapes.add_shape(RECT, Emu(0), Emu(0), W, Inches(0.11))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    if kicker:
        _text(slide, Inches(0.7), Inches(0.3), Inches(11.9), Inches(0.3),
              kicker.upper(), size=11, bold=True, colour=GOLD)
    _text(slide, Inches(0.7), Inches(0.56), Inches(11.9), Inches(0.7),
          title, size=27, bold=True, colour=NAVY)


def _card(slide, left, top, width, height, value: str, label: str,
          accent: RGBColor, value_size=34):
    shape = slide.shapes.add_shape(ROUND, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MIST
    shape.line.color.rgb = MIST
    shape.shadow.inherit = False
    stripe = slide.shapes.add_shape(RECT, left, top, Inches(0.06), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_top = Inches(0.22)
    head = frame.paragraphs[0]
    head.text = value
    head.alignment = PP_ALIGN.CENTER
    head.font.size = Pt(value_size)
    head.font.bold = True
    head.font.color.rgb = accent
    sub = frame.add_paragraph()
    sub.text = label
    sub.alignment = PP_ALIGN.CENTER
    sub.font.size = Pt(11.5)
    sub.font.bold = True
    sub.font.color.rgb = MUTED


def _picture(slide, image: Path, left, top, max_w, max_h) -> None:
    if not image or not Path(image).exists():
        return
    from PIL import Image

    with Image.open(image) as img:
        px_w, px_h = img.size
    scale = min(max_w / px_w, max_h / px_h)
    width, height = int(px_w * scale), int(px_h * scale)
    slide.shapes.add_picture(str(image),
                             left + int((max_w - width) / 2),
                             top + int((max_h - height) / 2),
                             width=width, height=height)


def _fmt(value, digits=1, suffix="") -> str:
    return "—" if value is None else f"{value:.{digits}f}{suffix}"


def _signed(value, digits=2) -> str:
    return "—" if value is None else f"{value:+.{digits}f}"


def _table(slide, rows: list[dict], columns: list[tuple[str, str]], top=Inches(1.55)):
    """A plain scoreboard. `columns` is [(heading, row-key)]."""
    table = slide.shapes.add_table(
        len(rows) + 1, len(columns), Inches(0.7), top,
        Inches(11.9), Inches(0.38) * (len(rows) + 1)
    ).table
    for col, (heading, _) in enumerate(columns):
        cell = table.cell(0, col)
        cell.text = heading
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11.5)
        para.font.bold = True
        para.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, start=1):
        for c, (_, key) in enumerate(columns):
            cell = table.cell(r, c)
            cell.text = str(row.get(key, "—"))
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10.5)
            para.font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else MIST
    return table


def generate_cohort_ppt(*, report: dict, generated_at: str = "") -> bytes:
    """Build the outcome and impact deck from a build_cohort_report() result."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    scope = report.get("scope", {})
    journey = report.get("journey", {})
    outcome = report.get("outcome", {})
    impact = report.get("impact", {})
    movement = report.get("movement", {})
    departments = report.get("departments", [])
    leaders = report.get("leaders", {})

    # ── Cover ────────────────────────────────────────────────────────────────
    slide = _blank(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    rule = slide.shapes.add_shape(RECT, Emu(0), Inches(2.55), W, Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = GOLD
    rule.line.fill.background()

    _text(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.4),
          "HACRI-E  ·  JAIN (DEEMED-TO-BE UNIVERSITY)", size=13, bold=True, colour=GOLD)
    _text(slide, Inches(0.9), Inches(1.72), Inches(11.5), Inches(0.9),
          "AI Literacy Outcome and Impact", size=40, bold=True, colour=WHITE)
    _text(slide, Inches(0.9), Inches(2.78), Inches(11.5), Inches(0.5),
          f"{scope.get('campus', '')}  ·  {scope.get('dept', '')}  ·  {scope.get('level', '')}",
          size=19, colour=WHITE)
    _text(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.4),
          f"{journey.get('registered', 0)} students registered  ·  "
          f"generated {generated_at or datetime.now():%d %B %Y}"
          if not generated_at else
          f"{journey.get('registered', 0)} students registered  ·  generated {generated_at}",
          size=13, italic=True, colour=RGBColor(0xB4, 0xC0, 0xD4))

    covers = [
        (str(journey.get("stages", [{}])[1].get("count", 0)) if journey.get("stages") else "0",
         "Baseline completed"),
        (str(journey.get("fully_done", 0)), "Post survey completed"),
        (_fmt(outcome.get("avg_overall"), 2), "Baseline score / 5"),
        (_signed((impact.get("avg_overall") - outcome.get("avg_overall"))
                 if (impact.get("avg_overall") is not None
                     and outcome.get("avg_overall") is not None) else None),
         "Change after workshop"),
    ]
    for i, (value, label) in enumerate(covers):
        left = Inches(0.9) + i * Inches(3.0)
        box = slide.shapes.add_shape(ROUND, left, Inches(4.25), Inches(2.75), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x14, 0x2C, 0x59)
        box.line.color.rgb = RGBColor(0x2A, 0x40, 0x69)
        box.shadow.inherit = False
        frame = box.text_frame
        frame.margin_top = Inches(0.4)
        head = frame.paragraphs[0]
        head.text = value
        head.alignment = PP_ALIGN.CENTER
        head.font.size = Pt(32)
        head.font.bold = True
        head.font.color.rgb = GOLD
        sub = frame.add_paragraph()
        sub.text = label
        sub.alignment = PP_ALIGN.CENTER
        sub.font.size = Pt(11.5)
        sub.font.color.rgb = WHITE

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)

        # ── The journey ──────────────────────────────────────────────────────
        slide = _blank(prs)
        _heading(slide, "How the cohort moved through the programme", "The journey")
        _picture(slide, plot_journey(journey, tmpdir / "journey.png"),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.4))
        for i, (value, label, tone) in enumerate([
            (str(journey.get("pending_pre", 0)), "Baseline still to complete", ROSE),
            (str(journey.get("pending_orientation", 0)), "Deeksharambh still to complete", GOLD),
            (str(journey.get("pending_post", 0)), "Post survey still to complete", DEEP),
            (_fmt(journey.get("completion_pct"), 0, "%"), "Finished all three steps", TEAL),
        ]):
            _card(slide, Inches(0.7) + i * Inches(3.15), Inches(5.15),
                  Inches(2.85), Inches(1.75), value, label, tone)

        # ── Campus split ─────────────────────────────────────────────────────
        campuses = [c for c in report.get("campuses", []) if c.get("registered")]
        if campuses:
            slide = _blank(prs)
            _heading(slide, "Bangalore and Kochi, step by step", "Campus split")
            _picture(slide, plot_campus_split(campuses, tmpdir / "campus.png", title=""),
                     Inches(0.7), Inches(1.5), Inches(11.9), Inches(3.5))
            _table(slide, [{
                "campus": clean(c["campus"], 26),
                "registered": c["registered"],
                "pre": c["pre"],
                "orientation": c["orientation"],
                "post": c["post"],
                "pending": c["pending_post"],
                "completion": _fmt(c["completion_pct"], 0, "%"),
            } for c in campuses], [
                ("Campus", "campus"), ("Registered", "registered"), ("Baseline", "pre"),
                ("Deeksharambh", "orientation"), ("Post", "post"),
                ("Post pending", "pending"), ("Completed all", "completion"),
            ], top=Inches(5.2))

        # ── Outcome ──────────────────────────────────────────────────────────
        slide = _blank(prs)
        _heading(slide, "Where the cohort started", "Outcome · baseline survey")
        for i, (value, label, tone) in enumerate([
            (str(outcome.get("scored", 0)), "Scored responses", NAVY),
            (_fmt(outcome.get("avg_lit"), 2), "AI Literacy / 5", DEEP),
            (_fmt(outcome.get("avg_read"), 2), "AI Readiness / 5", TEAL),
            (_fmt(outcome.get("avg_overall"), 2), "Overall / 5", GOLD),
        ]):
            _card(slide, Inches(0.7) + i * Inches(3.15), Inches(1.5),
                  Inches(2.85), Inches(1.6), value, label, tone)
        _picture(slide, plot_distribution(outcome, tmpdir / "outcome.png", colour="#1e3a8a"),
                 Inches(0.7), Inches(3.3), Inches(11.9), Inches(3.9))

        slide = _blank(prs)
        _heading(slide, "Who the baseline cohort were", "Outcome · quadrants and bands")
        _picture(slide, plot_mix(outcome, impact, tmpdir / "quadmix.png", "quadrants", title=""),
                 Inches(0.4), Inches(1.5), Inches(6.3), Inches(5.4))
        _picture(slide, plot_mix(outcome, impact, tmpdir / "bandmix.png", "bands", title=""),
                 Inches(6.7), Inches(1.5), Inches(6.3), Inches(5.4))

        # ── Impact ───────────────────────────────────────────────────────────
        slide = _blank(prs)
        _heading(slide, "Where the cohort ended up", "Impact · post-workshop survey")
        delta_overall = (impact.get("avg_overall") - outcome.get("avg_overall")
                         if (impact.get("avg_overall") is not None
                             and outcome.get("avg_overall") is not None) else None)
        for i, (value, label, tone) in enumerate([
            (str(impact.get("scored", 0)), "Scored responses", NAVY),
            (_fmt(impact.get("avg_lit"), 2), "AI Literacy / 5", DEEP),
            (_fmt(impact.get("avg_read"), 2), "AI Readiness / 5", TEAL),
            (_signed(delta_overall), "Overall change", GOLD),
        ]):
            _card(slide, Inches(0.7) + i * Inches(3.15), Inches(1.5),
                  Inches(2.85), Inches(1.6), value, label, tone)
        _picture(slide, plot_distribution(impact, tmpdir / "impact.png", colour="#0d9488"),
                 Inches(0.7), Inches(3.3), Inches(11.9), Inches(3.9))

        slide = _blank(prs)
        _heading(slide, "The shift, baseline against post", "Impact · the shape of the change")
        _picture(slide, plot_pre_post(outcome, impact, tmpdir / "prepost.png", title=""),
                 Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))

        slide = _blank(prs)
        _heading(slide, "How far each student moved", "Impact · matched students")
        _picture(slide, plot_movement(movement, tmpdir / "movement.png", title=""),
                 Inches(0.7), Inches(1.5), Inches(11.9), Inches(3.9))
        for i, (value, label, tone) in enumerate([
            (str(movement.get("matched", 0)), "Completed both surveys", NAVY),
            (_fmt(movement.get("gained_pct"), 0, "%"), "Improved", TEAL),
            (_signed(movement.get("avg_delta_lit")), "Average literacy change", DEEP),
            (str(movement.get("moved_quadrant", 0)), "Changed quadrant", GOLD),
        ]):
            _card(slide, Inches(0.7) + i * Inches(3.15), Inches(5.5),
                  Inches(2.85), Inches(1.5), value, label, tone)

        # ── Departments ──────────────────────────────────────────────────────
        if departments:
            slide = _blank(prs)
            _heading(slide, "Completion by department", "Departments · who filled")
            _picture(slide, plot_departments(departments, tmpdir / "deptdone.png",
                                             "completion_pct", title=""),
                     Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))

            movers = [d for d in departments if d["delta"] is not None]
            if movers:
                slide = _blank(prs)
                _heading(slide, "Score change by department", "Departments · impact")
                _picture(slide, plot_departments(movers, tmpdir / "deptdelta.png",
                                                 "delta", title=""),
                         Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))

            slide = _blank(prs)
            _heading(slide, "Department scoreboard", "Departments · every number")
            _table(slide, [{
                "dept": clean(d["dept"], 34),
                "registered": d["registered"],
                "pre": d["pre"],
                "orientation": d["orientation"],
                "post": d["post"],
                "pending": d["pending_post"],
                "completion": _fmt(d["completion_pct"], 0, "%"),
                "pre_avg": _fmt(d["pre_avg"], 2),
                "post_avg": _fmt(d["post_avg"], 2),
                "delta": _signed(d["delta"]),
            } for d in departments[:12]], [
                ("Department", "dept"), ("Registered", "registered"), ("Baseline", "pre"),
                ("Deeksharambh", "orientation"), ("Post", "post"), ("Post pending", "pending"),
                ("Completed", "completion"), ("Baseline / 5", "pre_avg"),
                ("Post / 5", "post_avg"), ("Change", "delta"),
            ])

        # ── Closing ──────────────────────────────────────────────────────────
        slide = _blank(prs)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = NAVY
        _text(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4),
              "IN SHORT", size=13, bold=True, colour=GOLD)
        lines = [
            f"{journey.get('registered', 0)} students registered; "
            f"{journey.get('fully_done', 0)} finished all three steps "
            f"({_fmt(journey.get('completion_pct'), 0, '%')}).",
            f"The baseline scored {_fmt(outcome.get('avg_overall'), 2)} out of 5; "
            f"after the workshop it was {_fmt(impact.get('avg_overall'), 2)} "
            f"({_signed(delta_overall)}).",
            f"{movement.get('gained', 0)} of {movement.get('matched', 0)} matched students improved.",
        ]
        if leaders.get("most_complete"):
            lines.append(
                f"{clean(leaders['most_complete']['dept'], 40)} is furthest along; "
                f"{clean((leaders.get('least_complete') or {}).get('dept', '—'), 40)} has the most left to do."
            )
        box = slide.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11.5), Inches(3.0))
        frame = box.text_frame
        frame.word_wrap = True
        for i, line in enumerate(lines):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = line
            para.font.size = Pt(20)
            para.font.bold = i == 0
            para.font.color.rgb = WHITE
            para.space_after = Pt(12)
        _text(slide, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.4),
              "Office of Academics · JAIN (Deemed-to-be University)",
              size=12.5, colour=RGBColor(0xB4, 0xC0, 0xD4))

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
