"""
orientation_count_export.py — the plain department headcount, as a deck and a
Word document.

Not the orientation report: no vibe, no NPS, no sessions. Just how many
answered, department by department, campus by campus — built from
`department_count_summary()` so this number can never drift from the report's
own. Its own palette too: a dark cover in the "Human & AI" house style
(near-black ground, one bold headline, a teal kicker line), carried through
every slide and page rather than left behind after the cover.
"""
from __future__ import annotations

import io
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# python-docx backs the Word half of this module only. Importing it at module
# load time meant a server without it installed couldn't even build the
# slide deck, since importing this file at all would fail — one missing
# optional dependency taking down an unrelated download. It loads lazily,
# the first time a docx is actually built.
try:
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches as DocxInches
    from docx.shared import Pt as DocxPt
    from docx.shared import RGBColor as DocxRGBColor
except ImportError:
    Document = WD_TABLE_ALIGNMENT = WD_ALIGN_PARAGRAPH = None
    OxmlElement = qn = DocxInches = DocxPt = None

    def DocxRGBColor(*_a, **_kw):  # placeholder so the module-level palette below still loads
        return None

# ── Palette ───────────────────────────────────────────────────────────────
INK = RGBColor(0x06, 0x08, 0x0F)        # the ground, corner to corner
PANEL = RGBColor(0x10, 0x15, 0x22)      # a card, one shade up from the ground
LINE = RGBColor(0x22, 0x2B, 0x3D)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)       # the one accent — kickers, rules, numbers
TEAL_DIM = RGBColor(0x17, 0x7F, 0x74)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9A, 0xA5, 0xB8)

DOCX_INK = DocxRGBColor(0x06, 0x08, 0x0F)
DOCX_TEAL = DocxRGBColor(0x2D, 0xD4, 0xBF)
DOCX_WHITE = DocxRGBColor(0xFF, 0xFF, 0xFF)
DOCX_INK_TEXT = DocxRGBColor(0x14, 0x18, 0x21)
DOCX_MUTED = DocxRGBColor(0x5B, 0x65, 0x70)

W = Inches(13.333)
H = Inches(7.5)
RECT = 1  # MSO_SHAPE.RECTANGLE


# ── PPTX ──────────────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    return slide


def _text(slide, left, top, width, height, text, *, size=14, bold=False,
          colour=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          spacing=0, italic=False, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        if line_spacing:
            para.line_spacing = line_spacing
        for font in [para.font] + [r.font for r in para.runs]:
            font.size = Pt(size)
            font.bold = bold
            font.italic = italic
            font.color.rgb = colour
        if spacing:
            # Letter-spacing has no python-pptx property; it lives in run XML.
            for run in para.runs or [para.add_run()]:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(spacing * 100)))
    return box


def _kicker(slide, left, top, text, *, width=Inches(11.5)):
    _text(slide, left, top, width, Inches(0.35), text.upper(),
          size=12, bold=True, colour=TEAL, spacing=2.2)


def _rule(slide, left, top, width, colour=TEAL, height=Emu(20000)):
    bar = slide.shapes.add_shape(RECT, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _panel(slide, left, top, width, height, *, fill=PANEL, line=LINE):
    box = slide.shapes.add_shape(RECT, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box


def _stat_card(slide, left, top, width, height, value: str, label: str):
    _panel(slide, left, top, width, height)
    _text(slide, left + Inches(0.18), top + Inches(0.16), width - Inches(0.36), Inches(0.7),
          value, size=30, bold=True, colour=TEAL)
    _text(slide, left + Inches(0.18), top + height - Inches(0.5), width - Inches(0.36), Inches(0.35),
          label.upper(), size=10.5, bold=True, colour=MUTED, spacing=0.8)


def _footer(slide, campus: str, generated_at: str, page: str = ""):
    _text(slide, Inches(0.7), H - Inches(0.5), Inches(9), Inches(0.35),
          f"Deeksharambh 2026 · {campus} · generated {generated_at}",
          size=9.5, colour=MUTED)
    if page:
        _text(slide, W - Inches(1.6), H - Inches(0.5), Inches(0.9), Inches(0.35),
              page, size=9.5, colour=MUTED, align=PP_ALIGN.RIGHT)


def _heading(slide, title: str, kicker: str):
    _kicker(slide, Inches(0.7), Inches(0.55), kicker)
    _text(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.7),
          title, size=30, bold=True, colour=WHITE)
    _rule(slide, Inches(0.7), Inches(1.55), Inches(0.9))


def _dept_table_slide(prs, campus_row: dict, generated_at: str, part: int, of: int):
    slide = _blank(prs)
    label = campus_row["campus"]
    suffix = f" — part {part} of {of}" if of > 1 else ""
    _heading(slide, f"{label}{suffix}", f"Department response count · {campus_row['count']} answered")

    rows = campus_row["rows"]
    n = len(rows) + 1
    table_shape = slide.shapes.add_table(
        n, 3, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.42) * n)
    table = table_shape.table
    table.columns[0].width = Inches(0.9)
    table.columns[1].width = Inches(8.7)
    table.columns[2].width = Inches(2.3)

    headers = ["#", "Department", "Responses"]
    for c, label_ in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = label_
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL_DIM
        if c == 2:
            para.alignment = PP_ALIGN.RIGHT

    for r, row in enumerate(rows, start=1):
        values = [str(row["rank"]), row["dept"], str(row["count"])]
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text = value
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11.5)
            para.font.color.rgb = WHITE if r % 2 else MUTED
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if r % 2 else PANEL
            if c == 2:
                para.alignment = PP_ALIGN.RIGHT

    _footer(slide, campus_row["campus"], generated_at)


def build_department_count_pptx(data: dict, *, generated_at: str) -> bytes:
    """The deck: a cover, one overview slide, then each campus's headcount."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ── Cover ────────────────────────────────────────────────────────────
    slide = _blank(prs)
    _rule(slide, Inches(0), Inches(0), W, height=Emu(46000))
    _kicker(slide, Inches(0.9), Inches(1.55), f"Deeksharambh 2026 · JAIN University · {data['campus']}")
    _text(slide, Inches(0.85), Inches(2.05), Inches(11.5), Inches(1.7),
          "DEPARTMENT\nRESPONSE COUNT", size=52, bold=True, colour=WHITE, spacing=0.4)
    _rule(slide, Inches(0.9), Inches(3.85), Inches(1.4))
    _text(slide, Inches(0.9), Inches(4.05), Inches(11), Inches(0.45),
          "EVERY DEPARTMENT · EVERY CAMPUS · NO SCORES, JUST WHO ANSWERED",
          size=14, bold=True, colour=TEAL, spacing=1.6)
    _text(slide, Inches(0.9), Inches(4.65), Inches(11.5), Inches(0.5),
          f"{data['total_answered']} of {data['total_registered']} students answered "
          f"({data['response_rate']:.0f}%) · generated {generated_at}",
          size=13, colour=MUTED)

    # ── Overview ─────────────────────────────────────────────────────────
    slide = _blank(prs)
    _heading(slide, "How many, where", "Overview")
    cards = [
        (str(data["total_registered"]), "Registered"),
        (str(data["total_answered"]), "Answered"),
        (str(data["total_pending"]), "Still pending"),
        (f"{data['response_rate']:.0f}%", "Response rate"),
        (str(data["ug"]), "UG"),
        (str(data["pg"]), "PG"),
    ]
    card_w = Inches(1.87)
    gap = Inches(0.14)
    for i, (value, label) in enumerate(cards):
        left = Inches(0.7) + i * (card_w + gap)
        _stat_card(slide, left, Inches(1.9), card_w, Inches(1.55), value, label)

    n = len(data["departments"])
    _text(slide, Inches(0.7), Inches(3.75), Inches(11.5), Inches(0.4),
          f"{n} department{'' if n == 1 else 's'} answered across "
          f"{len(data['campuses'])} campus{'' if len(data['campuses']) == 1 else 'es'}.",
          size=13, colour=MUTED)
    _footer(slide, data["campus"], generated_at)

    # ── Per campus ───────────────────────────────────────────────────────
    PAGE_SIZE = 18
    for row in data["campuses"]:
        depts = row["departments"]
        ranked = [{"rank": i + 1, **d} for i, d in enumerate(depts)]
        pages = [ranked[i:i + PAGE_SIZE] for i in range(0, len(ranked), PAGE_SIZE)] or [[]]
        for part, page_rows in enumerate(pages, start=1):
            _dept_table_slide(
                prs, {"campus": row["campus"], "count": row["count"], "rows": page_rows},
                generated_at, part, len(pages))

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ── DOCX ──────────────────────────────────────────────────────────────────

def _shade_cell(cell, hex_colour: str) -> None:
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_colour)
    cell._tc.get_or_add_tcPr().append(shd)


def _no_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "nil")
        borders.append(el)
    tbl_pr.append(borders)


def _run(paragraph, text, *, size=11, bold=False, colour=DOCX_INK_TEXT,
         align=None, spacing=None):
    """Add one styled run to an existing paragraph."""
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run(text)
    run.font.size = DocxPt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    if spacing:
        rPr = run._r.get_or_add_rPr()
        spc = OxmlElement("w:spacing")
        spc.set(qn("w:val"), str(int(spacing * 20)))
        rPr.append(spc)
    return run


def _cell_text(cell, text, **kw):
    """Clear a table cell and set its whole content to one styled run."""
    cell.text = ""
    return _run(cell.paragraphs[0], text, **kw)


def _cover_page(doc: Document, data: dict, generated_at: str) -> None:
    """A full-width dark cell, styled like the deck's own cover.

    Word has no reliable print-time page background, so the cover is content
    — one borderless table, one cell, shaded solid — which prints exactly as
    it looks, the same way a coloured table cell always does.
    """
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    _no_borders(table)
    cell = table.cell(0, 0)
    _shade_cell(cell, "06080F")
    cell.width = DocxInches(6.5)

    tcPr = cell._tc.get_or_add_tcPr()
    margins = OxmlElement("w:tcMar")
    for side, val in (("top", 720), ("bottom", 720), ("left", 460), ("right", 460)):
        el = OxmlElement(f"w:{side}")
        el.set(qn("w:w"), str(val))
        el.set(qn("w:type"), "dxa")
        margins.append(el)
    tcPr.append(margins)

    _run(cell.paragraphs[0], f"DEEKSHARAMBH 2026 · JAIN UNIVERSITY · {data['campus']}",
         size=10.5, bold=True, colour=DOCX_TEAL, spacing=1.2)

    p = cell.add_paragraph()
    p.paragraph_format.space_before = DocxPt(14)
    _run(p, "DEPARTMENT RESPONSE COUNT", size=28, bold=True, colour=DOCX_WHITE, spacing=0.6)

    p = cell.add_paragraph()
    p.paragraph_format.space_before = DocxPt(10)
    _run(p, "EVERY DEPARTMENT · EVERY CAMPUS · NO SCORES, JUST WHO ANSWERED",
         size=10.5, bold=True, colour=DOCX_TEAL, spacing=1)

    p = cell.add_paragraph()
    p.paragraph_format.space_before = DocxPt(12)
    _run(p, f"{data['total_answered']} of {data['total_registered']} students answered "
            f"({data['response_rate']:.0f}%) · generated {generated_at}",
         size=10, colour=DocxRGBColor(0x9A, 0xA5, 0xB8))

    doc.add_page_break()


def _overview_table(doc: Document, data: dict) -> None:
    h = doc.add_heading("How many, where", level=1)
    for run in h.runs:
        run.font.color.rgb = DOCX_INK
    doc.add_paragraph("Overview — every student in scope, counted once.")

    rows = [
        ("Registered", data["total_registered"]),
        ("Answered", data["total_answered"]),
        ("Still pending", data["total_pending"]),
        ("Response rate", f"{data['response_rate']:.0f}%"),
        ("UG", data["ug"]),
        ("PG", data["pg"]),
    ]
    table = doc.add_table(rows=len(rows), cols=2)
    table.style = "Light Grid Accent 1"
    for i, (label, value) in enumerate(rows):
        _cell_text(table.cell(i, 0), label, size=11, bold=True)
        _cell_text(table.cell(i, 1), str(value), size=11,
                   align=WD_ALIGN_PARAGRAPH.RIGHT)
    doc.add_paragraph()


def _campus_table(doc: Document, row: dict) -> None:
    h = doc.add_heading(f"{row['campus']} — department response count", level=1)
    for run in h.runs:
        run.font.color.rgb = DOCX_INK
    doc.add_paragraph(f"{row['count']} student{'s' if row['count'] != 1 else ''} answered, "
                       f"across {len(row['departments'])} department"
                       f"{'s' if len(row['departments']) != 1 else ''}.")

    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    widths = (DocxInches(0.7), DocxInches(4.6), DocxInches(1.4))
    header_labels = ("#", "Department", "Responses")
    for c, (label, width) in enumerate(zip(header_labels, widths)):
        cell = table.rows[0].cells[c]
        cell.width = width
        _shade_cell(cell, "177F74")
        align = WD_ALIGN_PARAGRAPH.RIGHT if c == 2 else None
        _cell_text(cell, label, size=11, bold=True, colour=DOCX_WHITE, align=align)

    for i, dept in enumerate(row["departments"], start=1):
        cells = table.add_row().cells
        for c, width in enumerate(widths):
            cells[c].width = width
        _cell_text(cells[0], str(i), size=10.5)
        _cell_text(cells[1], dept["dept"], size=10.5)
        _cell_text(cells[2], str(dept["count"]), size=10.5, align=WD_ALIGN_PARAGRAPH.RIGHT)
        if i % 2 == 0:
            for cell in cells:
                _shade_cell(cell, "F2F5F7")

    total_cells = table.add_row().cells
    _cell_text(total_cells[0], "", size=10.5)
    _cell_text(total_cells[1], "Total", size=10.5, bold=True)
    _cell_text(total_cells[2], str(row["count"]), size=10.5, bold=True,
               align=WD_ALIGN_PARAGRAPH.RIGHT)
    for cell in total_cells:
        _shade_cell(cell, "E4EEEC")
    doc.add_paragraph()


def build_department_count_docx(data: dict, *, generated_at: str) -> bytes:
    """The same headcount as the deck, laid out as a Word document."""
    if Document is None:
        raise RuntimeError(
            "python-docx is not installed — run `pip install -r requirements.txt`")
    doc = Document()
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    for margin in ("left_margin", "right_margin", "top_margin", "bottom_margin"):
        setattr(section, margin, DocxInches(0.9))

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)
    style.font.color.rgb = DOCX_INK_TEXT

    _cover_page(doc, data, generated_at)
    _overview_table(doc, data)
    for row in data["campuses"]:
        _campus_table(doc, row)

    footer = doc.sections[0].footer
    footer.paragraphs[0].text = f"Deeksharambh 2026 · {data['campus']} · generated {generated_at}"
    footer.paragraphs[0].runs[0].font.size = DocxPt(8.5)
    footer.paragraphs[0].runs[0].font.color.rgb = DOCX_MUTED

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()
